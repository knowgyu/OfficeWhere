from __future__ import annotations

import shutil
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parent.parent
OUTPUT_DIR = ROOT / "examples" / "demo_cases"
MANUAL_LIBRARY_DIR = ROOT / "examples" / "officewhere_test_library"


def build_excel_cases() -> None:
    versions = {
        1: {
            "headers": ["과제명", "담당자", "예산", "상태"],
            "rows": [
                ["AI 플랫폼 구축", "김철수", 320000000, "진행중"],
                ["데이터 허브 고도화", "이영희", 180000000, "완료"],
                ["보안 체계 개편", "박민수", 95000000, "검토중"],
            ],
            "far_cell": None,
        },
        2: {
            "headers": ["과제명", "담당자", "예산", "상태", "리스크"],
            "rows": [
                ["AI 플랫폼 구축", "김철수", 300000000, "진행중", "예산 조정"],
                ["데이터 허브 고도화", "이영희", 180000000, "완료", ""],
                ["차세대 포털 전환", "최은지", 210000000, "신규", "일정 확인 필요"],
            ],
            "far_cell": (11, 11, "v2 원거리 확인값"),
        },
        3: {
            "headers": ["과제명", "담당자", "예산", "상태", "리스크", "완료예정"],
            "rows": [
                ["AI 플랫폼 구축", "김철수", 310000000, "진행중", "벤더 계약 지연", "260510"],
                ["데이터 허브 고도화", "이영희", 185000000, "운영전환", "", "260503"],
                ["차세대 포털 전환", "최은지", 230000000, "진행중", "UX 검토 필요", "260531"],
                ["보안 체계 개편", "박민수", 110000000, "재개", "승인 대기", "260524"],
            ],
            "far_cell": (13, 12, "v3 원거리 확인값"),
        },
        4: {
            "headers": ["과제명", "담당자", "예산", "상태", "리스크", "완료예정", "비고"],
            "rows": [
                ["AI 플랫폼 구축", "김철수", 315000000, "검수중", "해소", "260517", "통합 테스트 진행"],
                ["차세대 포털 전환", "최은지", 240000000, "진행중", "디자인 확정 필요", "260607", "범위 확대"],
                ["보안 체계 개편", "박민수", 115000000, "진행중", "정책 검토", "260531", "재개 후 예산 증액"],
                ["운영 이관 준비", "한지민", 65000000, "신규", "인수인계 일정", "260614", "v4 신규 항목"],
            ],
            "far_cell": (15, 14, "v4 원거리 확인값"),
        },
    }

    for revision, data in versions.items():
        path = OUTPUT_DIR / f"excel_budget_v{revision}.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = "사업현황"

        ws["A1"] = "2026 상반기 사업 현황"
        ws["A2"] = "실제 표는 C3부터 시작한다."

        for offset, header in enumerate(data["headers"], start=3):
            ws.cell(row=3, column=offset, value=header)

        for row_index, row in enumerate(data["rows"], start=4):
            for col_index, value in enumerate(row, start=3):
                ws.cell(row=row_index, column=col_index, value=value)

        if data["far_cell"]:
            row, column, value = data["far_cell"]
            ws.cell(row=row, column=column, value=value)

        for column in ("A", "C", "D", "E", "F", "G", "H", "I", "L", "N"):
            ws.column_dimensions[column].width = 18
        wb.save(path)


def build_word_cases() -> None:
    documents = {
        1: {
            "paragraphs": [
                "본 문서는 1차 초안이다.",
                "예산 검토는 5월 말까지 완료한다.",
            ],
            "rows": [("항목", "내용"), ("담당", "플랫폼팀"), ("승인", "대기")],
        },
        2: {
            "paragraphs": [
                "본 문서는 2차 수정본이다.",
                "예산 검토는 6월 첫째 주까지 완료한다.",
                "보안 검토 항목이 추가되었다.",
            ],
            "rows": [("항목", "내용"), ("담당", "플랫폼팀"), ("승인", "완료"), ("보안", "추가 검토 필요")],
        },
        3: {
            "paragraphs": [
                "본 문서는 3차 검토본이다.",
                "예산 검토 결과를 반영했고 보안 검토 일정을 확정했다.",
                "운영 이관 체크리스트를 추가한다.",
            ],
            "rows": [("항목", "내용"), ("담당", "플랫폼팀"), ("승인", "조건부 승인"), ("보안", "검토 진행중"), ("운영", "체크리스트 작성")],
        },
        4: {
            "paragraphs": [
                "본 문서는 4차 릴리즈 후보이다.",
                "예산과 보안 검토가 완료되었고 운영 이관 준비를 시작한다.",
                "사용자 교육 일정이 추가되었다.",
            ],
            "rows": [("항목", "내용"), ("담당", "플랫폼팀"), ("승인", "완료"), ("보안", "완료"), ("운영", "이관 준비"), ("교육", "260524 예정")],
        },
    }

    for revision, data in documents.items():
        doc = Document()
        doc.add_heading("클라우드 전환 검토", level=1)
        for paragraph in data["paragraphs"]:
            doc.add_paragraph(paragraph)
        table = doc.add_table(rows=len(data["rows"]), cols=2)
        for row_index, (label, value) in enumerate(data["rows"]):
            table.cell(row_index, 0).text = label
            table.cell(row_index, 1).text = value
        doc.save(OUTPUT_DIR / f"proposal_note_v{revision}.docx")


def _add_slide(prs: Presentation, title: str, body_lines: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    for index, line in enumerate(body_lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line


def build_ppt_cases() -> None:
    decks = {
        1: [
            ("프로젝트 개요", ["범위 정의 완료", "예산 검토 진행중"]),
            ("주요 일정", ["5월: 요구사항 확정", "6월: 구현 시작"]),
        ],
        2: [
            ("프로젝트 개요", ["범위 정의 완료", "예산 검토 완료"]),
            ("위험 요소", ["인력 확보 지연", "외부 연동 일정 미정"]),
            ("주요 일정", ["5월: 요구사항 확정", "6월: 구현 시작", "7월: 사용자 테스트"]),
        ],
        3: [
            ("프로젝트 개요", ["검색/정합성 1차 사용자 피드백 반영", "Excel 표 보기 개선 진행"]),
            ("위험 요소", ["네트워크 드라이브 검증 필요", "Excel 대용량 표 표시 정책 검토"]),
            ("주요 일정", ["260503: v3 내부 검수", "260510: 사용자 피드백 수집"]),
        ],
        4: [
            ("프로젝트 개요", ["버전관리 표 보기 개선 완료", "릴리즈 후보 준비"]),
            ("주요 변경점", ["Excel v4 표 보기 추가 검증", "닫기/트레이 동작 정리"]),
            ("위험 요소", ["대용량 Excel 렌더링 성능 모니터링"]),
            ("주요 일정", ["260517: v4 검수", "260524: 사용자 배포 준비"]),
        ],
    }

    for revision, slides in decks.items():
        prs = Presentation()
        for title, lines in slides:
            _add_slide(prs, title, lines)
        prs.save(OUTPUT_DIR / f"status_review_v{revision}.pptx")

def _save_word_document(
    path: Path,
    heading: str,
    paragraphs: list[str],
    table_rows: list[tuple[str, str]] | None = None,
) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    doc = Document()
    doc.add_heading(heading, level=1)
    for paragraph in paragraphs:
        doc.add_paragraph(paragraph)
    if table_rows:
        table = doc.add_table(rows=len(table_rows), cols=2)
        for row_index, (label, value) in enumerate(table_rows):
            table.cell(row_index, 0).text = label
            table.cell(row_index, 1).text = value
    doc.save(path)


def _save_budget_workbook(path: Path, revision: int) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    ws = wb.active
    ws.title = "예산현황"

    ws["A1"] = "2026 문서 검색/정합성 고도화 예산"
    ws["A2"] = "앱의 엑셀 표 위치 자동 감지를 확인하기 위해 표는 C4부터 시작한다."

    versions: dict[int, dict[str, object]] = {
        1: {
            "headers": ["과제명", "담당자", "예산", "상태"],
            "rows": [
                ["A 프로젝트 검색 개선", "김철수", 120000000, "진행중"],
                ["Office 정합성 검사", "이영희", 85000000, "검토중"],
                ["배포 자동화", "박민수", 45000000, "예정"],
            ],
        },
        2: {
            "headers": ["과제명", "담당자", "예산", "상태", "변경사유"],
            "rows": [
                ["A 프로젝트 검색 개선", "김철수", 135000000, "진행중", "예산 조정"],
                ["Office 정합성 검사", "이영희", 85000000, "완료", "검증 완료"],
                ["사용자 교육 자료", "최은지", 30000000, "신규", "사용자 요청 추가"],
            ],
        },
        3: {
            "headers": ["과제명", "담당자", "예산", "상태", "변경사유", "완료예정"],
            "rows": [
                ["A 프로젝트 검색 개선", "김철수", 145000000, "진행중", "검색 범위 확대", "260503"],
                ["Office 정합성 검사", "이영희", 90000000, "완료", "사용자 피드백 반영", "260503"],
                ["사용자 교육 자료", "최은지", 35000000, "진행중", "사용자 요청 추가", "260510"],
                ["배포 자동화", "박민수", 50000000, "재개", "릴리즈 준비", "260517"],
            ],
        },
        4: {
            "headers": ["과제명", "담당자", "예산", "상태", "변경사유", "완료예정"],
            "rows": [
                ["A 프로젝트 검색 개선", "김철수", 150000000, "검수중", "성능 검증 추가", "260510"],
                ["Office 정합성 검사", "이영희", 92000000, "개선중", "버전관리 표 보기 개선", "260517"],
                ["사용자 교육 자료", "최은지", 42000000, "진행중", "교육 범위 확대", "260524"],
                ["네트워크 드라이브 검증", "한지민", 28000000, "신규", "사내 공유 드라이브 확인", "260531"],
            ],
        },
        5: {
            "headers": ["과제명", "담당자", "예산", "상태", "변경사유", "완료예정"],
            "rows": [
                ["A 프로젝트 검색 개선", "김철수", 155000000, "완료", "릴리즈 후보 반영", "260517"],
                ["Office 정합성 검사", "이영희", 98000000, "검수중", "Excel 표로 보기 개선", "260524"],
                ["네트워크 드라이브 검증", "한지민", 32000000, "진행중", "읽기 전용 안전성 확인", "260607"],
                ["운영 이관 준비", "오세훈", 60000000, "신규", "사용자 배포 준비", "260614"],
            ],
            "far_cell": (12, 10, "v4 확인용 원거리 셀"),
        },
    }
    data = versions[revision]

    for offset, header in enumerate(data["headers"], start=3):
        ws.cell(row=4, column=offset, value=header)

    for row_index, row in enumerate(data["rows"], start=5):
        for col_index, value in enumerate(row, start=3):
            ws.cell(row=row_index, column=col_index, value=value)

    if data.get("far_cell"):
        row, column, value = data["far_cell"]
        ws.cell(row=row, column=column, value=value)

    for column in ("A", "C", "D", "E", "F", "G", "H", "J"):
        ws.column_dimensions[column].width = 18
    wb.save(path)


def _save_common_form(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    ws = wb.active
    ws.title = "공통양식"
    ws.append(["문서ID", "제목", "담당자", "확인상태"])
    ws.append(["FORM-001", "공통양식 배포 확인", "총무팀", "동일본 확인용"])
    ws.append(["FORM-002", "A 프로젝트 교육 참석자", "교육팀", "확인 완료"])
    for column in ("A", "B", "C", "D"):
        ws.column_dimensions[column].width = 20
    wb.save(path)


def _save_project_status_deck(path: Path, revision: int) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    decks = {
        1: [
            ("A 프로젝트 상태", ["검색 범위 옵션 설계", "정합성 그룹 기준 검토"]),
            ("다음 일정", ["260419 기준: 내부 테스트 준비", "위험: 엑셀 예산 확정 전"]),
        ],
        2: [
            ("A 프로젝트 상태", ["본문만 검색 배포 완료", "Office 정합성 그룹 기준 적용"]),
            ("위험 요소", ["일정 지연 가능성 감소", "사용자 교육 필요"]),
            ("다음 일정", ["260426 기준: 사용자 검증", "릴리즈 체크리스트 완료"]),
        ],
        3: [
            ("A 프로젝트 상태", ["검색/정합성 1차 사용자 피드백 반영", "Excel 표 보기 개선 진행"]),
            ("위험 요소", ["공유 드라이브 접근성 확인", "대용량 Excel 표시 정책 검토"]),
            ("다음 일정", ["260503 기준: v2 내부 검수", "260510 기준: 사용자 피드백 수집"]),
        ],
        4: [
            ("A 프로젝트 상태", ["Excel 표로 보기 변경 이력 개선", "창 닫기/트레이 동작 안정화"]),
            ("주요 변경점", ["추가/삭제 셀 색상 표시", "상세 이력 수정 전/후 박스 표시"]),
            ("다음 일정", ["260510 기준: v3 검수", "260517 기준: 릴리즈 후보 준비"]),
        ],
        5: [
            ("A 프로젝트 상태", ["버전관리 표 보기 개선 완료", "릴리즈 후보 준비"]),
            ("주요 변경점", ["Excel v4 표 보기 추가 검증", "설정 화면 문구 정리 예정"]),
            ("위험 요소", ["대용량 Excel 렌더링 성능 모니터링"]),
            ("다음 일정", ["260517 기준: v4 검수", "260524 기준: 사용자 배포 준비"]),
        ],
    }
    for title, lines in decks[revision]:
        _add_slide(prs, title, lines)
    prs.save(path)


def build_manual_test_library() -> None:
    """Build a realistic folder tree that can be registered in the app."""
    if MANUAL_LIBRARY_DIR.exists():
        shutil.rmtree(MANUAL_LIBRARY_DIR)
    MANUAL_LIBRARY_DIR.mkdir(parents=True)

    _save_word_document(
        MANUAL_LIBRARY_DIR / "01_보고서" / "주간보고_v1.0_260419.docx",
        "주간보고",
        [
            "260419 기준 A 프로젝트 검색 개선은 설계 단계이다.",
            "본문만 검색과 파일명 검색을 분리하는 방안을 검토한다.",
            "보안 검토는 다음 주에 시작한다.",
        ],
        [("작성자", "운영팀"), ("상태", "초안")],
    )
    _save_word_document(
        MANUAL_LIBRARY_DIR / "01_보고서" / "주간보고_v1.1_260426.docx",
        "주간보고",
        [
            "260426 기준 A 프로젝트 검색 개선은 구현 완료 상태이다.",
            "본문만 검색으로 예산 조정, 보안 검토 같은 본문 키워드를 찾을 수 있다.",
            "보안 검토 항목이 추가되었고 사용자 테스트를 시작한다.",
        ],
        [("작성자", "운영팀"), ("상태", "검토 완료"), ("추가", "보안 검토 포함")],
    )
    _save_word_document(
        MANUAL_LIBRARY_DIR / "01_보고서" / "주간보고_v2.0_260503.docx",
        "주간보고",
        [
            "260503 기준 사용자 테스트 피드백을 반영했다.",
            "Excel 표로 보기와 공유 드라이브 읽기 정책을 검토한다.",
            "버전관리 화면의 보기 크기 설정을 추가했다.",
        ],
        [("작성자", "운영팀"), ("상태", "2차 검토"), ("추가", "Excel 표로 보기 검토")],
    )
    _save_word_document(
        MANUAL_LIBRARY_DIR / "01_보고서" / "주간보고_v3.0_260510.docx",
        "주간보고",
        [
            "260510 기준 Excel 표로 보기 개선을 완료했다.",
            "삭제/추가 셀 표시와 가로 스크롤 동작을 확인한다.",
            "사용자 검수에서 버전 이력 설명을 단순화한다.",
        ],
        [("작성자", "운영팀"), ("상태", "3차 검토"), ("추가", "표 보기 검수")],
    )
    _save_word_document(
        MANUAL_LIBRARY_DIR / "01_보고서" / "주간보고_v4.0_260517.docx",
        "주간보고",
        [
            "260517 기준 릴리즈 후보를 준비한다.",
            "검색과 정합성 검사의 주요 시나리오를 다시 점검한다.",
            "사용자 배포 전 체크리스트를 확정한다.",
        ],
        [("작성자", "운영팀"), ("상태", "릴리즈 후보"), ("추가", "사용자 배포 준비")],
    )

    _save_project_status_deck(
        MANUAL_LIBRARY_DIR / "01_보고서" / "프로젝트상태_v1.0.pptx",
        revision=1,
    )
    _save_project_status_deck(
        MANUAL_LIBRARY_DIR / "01_보고서" / "프로젝트상태_v1.1_260426.pptx",
        revision=2,
    )
    _save_project_status_deck(
        MANUAL_LIBRARY_DIR / "01_보고서" / "프로젝트상태_v2.0_260503.pptx",
        revision=3,
    )
    _save_project_status_deck(
        MANUAL_LIBRARY_DIR / "01_보고서" / "프로젝트상태_v3.0_260510.pptx",
        revision=4,
    )
    _save_project_status_deck(
        MANUAL_LIBRARY_DIR / "01_보고서" / "프로젝트상태_v4.0_260517.pptx",
        revision=5,
    )

    _save_budget_workbook(
        MANUAL_LIBRARY_DIR / "02_예산" / "사업예산_v1.0_260419.xlsx",
        revision=1,
    )
    _save_budget_workbook(
        MANUAL_LIBRARY_DIR / "02_예산" / "사업예산_v1.1_260426.xlsx",
        revision=2,
    )
    _save_budget_workbook(
        MANUAL_LIBRARY_DIR / "02_예산" / "사업예산_v2.0_260503.xlsx",
        revision=3,
    )
    _save_budget_workbook(
        MANUAL_LIBRARY_DIR / "02_예산" / "사업예산_v3.0_260510.xlsx",
        revision=4,
    )
    _save_budget_workbook(
        MANUAL_LIBRARY_DIR / "02_예산" / "사업예산_v4.0_260517.xlsx",
        revision=5,
    )

    _save_word_document(
        MANUAL_LIBRARY_DIR / "03_부서A" / "회의록.docx",
        "회의록",
        [
            "부서A 회의록: 안내 자료 검토와 A 프로젝트 데모 준비가 주요 안건이다.",
            "다음 주까지 사용자 안내 문구를 확정한다.",
        ],
        [("부서", "A"), ("결론", "데모 준비")],
    )
    _save_word_document(
        MANUAL_LIBRARY_DIR / "04_부서B" / "회의록.docx",
        "회의록",
        [
            "부서B 회의록: 개발 일정 지연과 예산 조정 요청이 주요 안건이다.",
            "260426 기준으로 리스크 보고서를 갱신한다.",
        ],
        [("부서", "B"), ("결론", "예산 조정")],
    )
    _save_common_form(MANUAL_LIBRARY_DIR / "03_부서A" / "공통양식.xlsx")
    _save_common_form(MANUAL_LIBRARY_DIR / "04_부서B" / "공통양식.xlsx")

    search_dir = MANUAL_LIBRARY_DIR / "05_검색샘플"
    search_dir.mkdir(parents=True, exist_ok=True)
    _save_word_document(
        search_dir / "운영메모.docx",
        "A 프로젝트 체험 메모",
        [
            "본문만 검색에서 이 문장이 검색되어야 한다.",
            "파일명에는 A 프로젝트가 없으므로 파일명만 검색에서는 잡히지 않는 것이 정상이다.",
        ],
    )
    _save_word_document(
        search_dir / "릴리즈체크리스트.docx",
        "릴리즈 체크리스트",
        [
            "본문만 검색: 보안 검토, 예산 조정, A 프로젝트 확인",
            "정합성 검사: 같은 파일명/다른 내용, 버전 이력 그룹 확인",
        ],
    )
    _save_word_document(
        MANUAL_LIBRARY_DIR / "README_테스트방법.docx",
        "OfficeWhere 실제 사례 테스트 라이브러리",
        [
            "앱에서 이 폴더를 문서 폴더로 추가한 뒤 문서 새로고침을 실행하세요.",
            "검색 확인 키워드: A 프로젝트(Word/PPT 본문), 예산 조정(Excel/Word 본문), 회의록(파일명)",
            "정합성 확인: 회의록.docx, 공통양식.xlsx, 주간보고/프로젝트상태/사업예산 버전 이력",
        ],
    )


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    build_excel_cases()
    build_word_cases()
    build_ppt_cases()
    build_manual_test_library()
    print(f"demo cases written to {OUTPUT_DIR}")
    print(f"manual test library written to {MANUAL_LIBRARY_DIR}")


if __name__ == "__main__":
    main()
