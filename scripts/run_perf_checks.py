from __future__ import annotations

import tempfile
from pathlib import Path
from time import perf_counter
import sys

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from backend.core.checker import run_consistency_check


def build_excel_pair(root: Path, rows: int = 1000) -> tuple[Path, Path, dict]:
    parser_config = {
        "sheet_name": "Data",
        "header_row": 3,
        "start_col": 3,
        "end_col": 6,
        "end_row": rows + 3,
    }
    headers = ["과제명", "담당자", "예산", "상태"]

    for path, budget_offset in ((root / "perf_a.xlsx", 0), (root / "perf_b.xlsx", 5)):
        workbook = Workbook()
        worksheet = workbook.active
        worksheet.title = "Data"
        worksheet["A1"] = "성능 검증용 시트"
        worksheet["A2"] = "표는 C3부터 시작"
        for column_index, header in enumerate(headers, start=3):
            worksheet.cell(row=3, column=column_index, value=header)
        for row_index in range(rows):
            worksheet.cell(row=row_index + 4, column=3, value=f"과제-{row_index:04d}")
            worksheet.cell(row=row_index + 4, column=4, value=f"담당-{row_index % 20:02d}")
            worksheet.cell(row=row_index + 4, column=5, value=1000 + row_index + budget_offset)
            worksheet.cell(row=row_index + 4, column=6, value="진행중" if row_index % 2 == 0 else "완료")
        workbook.save(path)
    return root / "perf_a.xlsx", root / "perf_b.xlsx", parser_config


def build_word_pair(root: Path, blocks: int = 300) -> tuple[Path, Path]:
    left = root / "perf_a.docx"
    right = root / "perf_b.docx"
    for path, changed in ((left, False), (right, True)):
        document = Document()
        for index in range(blocks):
            text = f"문단 {index:03d} 내용"
            if changed and index % 25 == 0:
                text += " 수정"
            document.add_paragraph(text)
        document.save(path)
    return left, right


def build_ppt_pair(root: Path, slides: int = 30) -> tuple[Path, Path]:
    def add_textbox(slide, text: str) -> None:
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(6), Inches(1))
        textbox.text_frame.text = text

    left = root / "perf_a.pptx"
    right = root / "perf_b.pptx"
    for path, changed in ((left, False), (right, True)):
        presentation = Presentation()
        layout = presentation.slide_layouts[5]
        for index in range(slides):
            slide = presentation.slides.add_slide(layout)
            slide.shapes.title.text = f"Slide {index:02d}"
            text = f"본문 {index:02d}"
            if changed and index % 7 == 0:
                text += " 변경"
            add_textbox(slide, text)
        if changed:
            slide = presentation.slides.add_slide(layout)
            slide.shapes.title.text = "Inserted"
            add_textbox(slide, "추가 슬라이드")
        presentation.save(path)
    return left, right


def benchmark_excel(root: Path) -> float:
    left, right, parser_config = build_excel_pair(root)
    started = perf_counter()
    run_consistency_check(
        [
            {"id": 1, "path": str(left), "name": left.name, "file_type": "Excel", "key_column": "과제명", "parser_config": parser_config},
            {"id": 2, "path": str(right), "name": right.name, "file_type": "Excel", "key_column": "과제명", "parser_config": parser_config},
        ]
    )
    return perf_counter() - started


def benchmark_word(root: Path) -> float:
    left, right = build_word_pair(root)
    started = perf_counter()
    run_consistency_check(
        [
            {"id": 1, "path": str(left), "name": left.name, "file_type": "Word", "key_column": "", "parser_config": {}},
            {"id": 2, "path": str(right), "name": right.name, "file_type": "Word", "key_column": "", "parser_config": {}},
        ]
    )
    return perf_counter() - started


def benchmark_ppt(root: Path) -> float:
    left, right = build_ppt_pair(root)
    started = perf_counter()
    run_consistency_check(
        [
            {"id": 1, "path": str(left), "name": left.name, "file_type": "PowerPoint", "key_column": "", "parser_config": {}},
            {"id": 2, "path": str(right), "name": right.name, "file_type": "PowerPoint", "key_column": "", "parser_config": {}},
        ]
    )
    return perf_counter() - started


def main() -> None:
    with tempfile.TemporaryDirectory() as tmp_dir:
        root = Path(tmp_dir)
        print(f"[perf] excel 1000 rows: {benchmark_excel(root):.3f}s")
        print(f"[perf] word 300 blocks: {benchmark_word(root):.3f}s")
        print(f"[perf] ppt 30 slides: {benchmark_ppt(root):.3f}s")


if __name__ == "__main__":
    main()
