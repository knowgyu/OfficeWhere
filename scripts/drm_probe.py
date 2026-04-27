#!/usr/bin/env python3
"""
Probe whether the current Python process can read DRM-protected OOXML files.

This intentionally uses only the Python standard library by default.  If this
script can open a .docx/.xlsx/.pptx as a zip and read representative XML parts,
the important DRM question is answered: this python.exe process is receiving the
decrypted Office document bytes.

Optional library checks are available with --library-check, but they are not
required for the DRM trust/proces-host proof.
"""

from __future__ import annotations

import argparse
import json
import platform
import re
import sys
import time
import traceback
import zipfile
from pathlib import Path
from typing import Any, Iterable
from xml.etree import ElementTree

TEXT_RE = re.compile(r"\s+")


def normalize_text(value: str) -> str:
    return TEXT_RE.sub(" ", value).strip()


def local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def collect_xml_text(xml_bytes: bytes, limit: int) -> tuple[int, list[str]]:
    try:
        root = ElementTree.fromstring(xml_bytes)
    except ElementTree.ParseError:
        return 0, []

    count = 0
    samples: list[str] = []
    for elem in root.iter():
        if local_name(elem.tag) != "t" or not elem.text:
            continue
        text = normalize_text(elem.text)
        if not text:
            continue
        count += 1
        if len(samples) < limit:
            samples.append(text[:160])
    return count, samples


def read_part_texts(archive: zipfile.ZipFile, part_names: Iterable[str], sample_limit: int) -> dict[str, Any]:
    parts: list[dict[str, Any]] = []
    total_text_nodes = 0
    samples: list[str] = []

    for name in part_names:
        try:
            text_count, part_samples = collect_xml_text(archive.read(name), max(0, sample_limit - len(samples)))
        except KeyError:
            continue
        total_text_nodes += text_count
        samples.extend(part_samples)
        parts.append({"name": name, "text_nodes": text_count})

    return {
        "parts_read": len(parts),
        "text_nodes": total_text_nodes,
        "parts": parts[:20],
        "samples": samples[:sample_limit],
    }


def inspect_docx(archive: zipfile.ZipFile, sample_limit: int) -> dict[str, Any]:
    names = set(archive.namelist())
    required = "word/document.xml"
    if required not in names:
        raise RuntimeError("missing word/document.xml; not a normal DOCX payload")
    parts = [required]
    parts.extend(sorted(name for name in names if name.startswith("word/header") and name.endswith(".xml")))
    parts.extend(sorted(name for name in names if name.startswith("word/footer") and name.endswith(".xml")))
    result = read_part_texts(archive, parts, sample_limit)
    result.update({"document_part_present": True})
    return result


def inspect_xlsx(archive: zipfile.ZipFile, sample_limit: int) -> dict[str, Any]:
    names = set(archive.namelist())
    required = "xl/workbook.xml"
    if required not in names:
        raise RuntimeError("missing xl/workbook.xml; not a normal XLSX/XLSM payload")

    workbook_text_count, workbook_samples = collect_xml_text(archive.read(required), sample_limit)
    shared_text_count = 0
    shared_samples: list[str] = []
    if "xl/sharedStrings.xml" in names:
        shared_text_count, shared_samples = collect_xml_text(
            archive.read("xl/sharedStrings.xml"), max(0, sample_limit - len(workbook_samples))
        )

    worksheets = sorted(name for name in names if re.fullmatch(r"xl/worksheets/sheet\d+\.xml", name))
    worksheet_parts = read_part_texts(archive, worksheets[:10], max(0, sample_limit - len(workbook_samples) - len(shared_samples)))

    return {
        "workbook_part_present": True,
        "worksheet_count_detected": len(worksheets),
        "workbook_text_nodes": workbook_text_count,
        "shared_string_text_nodes": shared_text_count,
        "worksheet_text_nodes_sampled": worksheet_parts["text_nodes"],
        "parts_read": 1 + (1 if "xl/sharedStrings.xml" in names else 0) + worksheet_parts["parts_read"],
        "parts": [{"name": required, "text_nodes": workbook_text_count}]
        + ([{"name": "xl/sharedStrings.xml", "text_nodes": shared_text_count}] if "xl/sharedStrings.xml" in names else [])
        + worksheet_parts["parts"],
        "samples": (workbook_samples + shared_samples + worksheet_parts["samples"])[:sample_limit],
    }


def inspect_pptx(archive: zipfile.ZipFile, sample_limit: int) -> dict[str, Any]:
    names = set(archive.namelist())
    required = "ppt/presentation.xml"
    if required not in names:
        raise RuntimeError("missing ppt/presentation.xml; not a normal PPTX payload")
    slides = sorted(
        (name for name in names if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)),
        key=lambda value: int(re.search(r"slide(\d+)\.xml", value).group(1)),
    )
    result = read_part_texts(archive, slides[:20], sample_limit)
    result.update({"presentation_part_present": True, "slide_count_detected": len(slides)})
    return result


def inspect_ooxml(path: Path, sample_limit: int) -> dict[str, Any]:
    suffix = path.suffix.lower()
    with zipfile.ZipFile(path) as archive:
        names = archive.namelist()
        base: dict[str, Any] = {
            "zip_opened": True,
            "zip_entry_count": len(names),
            "content_types_present": "[Content_Types].xml" in names,
        }
        if suffix == ".docx":
            base.update(inspect_docx(archive, sample_limit))
        elif suffix in {".xlsx", ".xlsm"}:
            base.update(inspect_xlsx(archive, sample_limit))
        elif suffix == ".pptx":
            base.update(inspect_pptx(archive, sample_limit))
        else:
            raise RuntimeError(f"unsupported extension for OOXML probe: {suffix}")
        return base


def optional_library_check(path: Path, sample_limit: int) -> dict[str, Any]:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        from docx import Document  # type: ignore

        document = Document(str(path))
        samples = [normalize_text(paragraph.text) for paragraph in document.paragraphs if normalize_text(paragraph.text)]
        return {"library": "python-docx", "paragraph_count": len(document.paragraphs), "samples": samples[:sample_limit]}

    if suffix in {".xlsx", ".xlsm"}:
        from openpyxl import load_workbook  # type: ignore

        workbook = load_workbook(str(path), read_only=True, data_only=True)
        try:
            sheet_names = list(workbook.sheetnames)
            sample_values: list[str] = []
            for sheet in workbook.worksheets[:3]:
                for row in sheet.iter_rows(max_row=20, values_only=True):
                    for value in row:
                        if value is None:
                            continue
                        sample_values.append(str(value)[:160])
                        if len(sample_values) >= sample_limit:
                            break
                    if len(sample_values) >= sample_limit:
                        break
                if len(sample_values) >= sample_limit:
                    break
            return {"library": "openpyxl", "sheet_names": sheet_names, "samples": sample_values}
        finally:
            workbook.close()

    if suffix == ".pptx":
        from pptx import Presentation  # type: ignore

        presentation = Presentation(str(path))
        samples: list[str] = []
        text_shape_count = 0
        for slide in presentation.slides:
            for shape in slide.shapes:
                text = normalize_text(getattr(shape, "text", "") or "")
                if not text:
                    continue
                text_shape_count += 1
                if len(samples) < sample_limit:
                    samples.append(text[:160])
        return {"library": "python-pptx", "slide_count": len(presentation.slides), "text_shape_count": text_shape_count, "samples": samples}

    raise RuntimeError(f"unsupported extension for library check: {suffix}")


def probe_file(path: Path, sample_limit: int, include_samples: bool, library_check: bool, include_traceback: bool) -> dict[str, Any]:
    started = time.perf_counter()
    result: dict[str, Any] = {
        "path": str(path),
        "suffix": path.suffix.lower(),
        "exists": path.exists(),
        "status": "unknown",
    }

    try:
        stat = path.stat()
        result.update({"size_bytes": stat.st_size, "mtime": stat.st_mtime})
        result["zipfile_is_zipfile"] = zipfile.is_zipfile(path)
        result["stdlib_ooxml"] = inspect_ooxml(path, sample_limit if include_samples else 0)
        if not include_samples:
            result["stdlib_ooxml"].pop("samples", None)
        if library_check:
            library = optional_library_check(path, sample_limit if include_samples else 0)
            if not include_samples:
                library.pop("samples", None)
            result["library_check"] = library
        result["status"] = "ok"
    except Exception as exc:  # noqa: BLE001 - diagnostic script should report all failures.
        result["status"] = "error"
        result["error_type"] = type(exc).__name__
        result["error"] = str(exc)
        if include_traceback:
            result["traceback"] = traceback.format_exc()
    finally:
        result["duration_ms"] = round((time.perf_counter() - started) * 1000, 1)
    return result


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Check whether this python.exe process can read DRM-protected Office OOXML files."
    )
    parser.add_argument("files", nargs="+", help="DRM-protected .docx/.xlsx/.xlsm/.pptx files to probe")
    parser.add_argument("--include-samples", action="store_true", help="Include short extracted text samples in JSON output")
    parser.add_argument("--library-check", action="store_true", help="Also try python-docx/openpyxl/python-pptx if installed")
    parser.add_argument("--traceback", action="store_true", help="Include Python traceback for failures")
    parser.add_argument("--json-out", help="Write the JSON report to this path as well as stdout")
    parser.add_argument("--sample-limit", type=int, default=5, help="Maximum text samples to include when --include-samples is set")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    files = [Path(raw).expanduser() for raw in args.files]
    report = {
        "probe": "officewhere-drm-python-probe",
        "python": {
            "executable": sys.executable,
            "version": sys.version,
            "implementation": platform.python_implementation(),
            "prefix": sys.prefix,
            "base_prefix": getattr(sys, "base_prefix", sys.prefix),
            "platform": platform.platform(),
        },
        "checks": {
            "stdlib_ooxml": True,
            "library_check": bool(args.library_check),
            "include_samples": bool(args.include_samples),
        },
        "files": [
            probe_file(
                path=path,
                sample_limit=max(0, args.sample_limit),
                include_samples=bool(args.include_samples),
                library_check=bool(args.library_check),
                include_traceback=bool(args.traceback),
            )
            for path in files
        ],
    }
    report["summary"] = {
        "total": len(report["files"]),
        "ok": sum(1 for item in report["files"] if item.get("status") == "ok"),
        "error": sum(1 for item in report["files"] if item.get("status") != "ok"),
    }

    output = json.dumps(report, ensure_ascii=False, indent=2)
    print(output)
    if args.json_out:
        out_path = Path(args.json_out).expanduser()
        out_path.parent.mkdir(parents=True, exist_ok=True)
        out_path.write_text(output + "\n", encoding="utf-8")

    return 0 if report["summary"]["error"] == 0 else 1


if __name__ == "__main__":
    raise SystemExit(main())
