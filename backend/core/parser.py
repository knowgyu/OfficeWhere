import os
from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional
import pandas as pd


SUPPORTED_EXTENSIONS = {".xlsx", ".xls", ".docx", ".pptx"}


def get_file_type(path: str) -> str:
    ext = Path(path).suffix.lower()
    mapping = {
        ".xlsx": "Excel",
        ".xls": "Excel",
        ".docx": "Word",
        ".pptx": "PowerPoint",
    }
    return mapping.get(ext, "Unknown")


def parse_excel(path: str) -> pd.DataFrame:
    """Excel 파일 파싱"""
    try:
        ext = Path(path).suffix.lower()
        if ext == ".xls":
            df = pd.read_excel(path, engine="xlrd")
        else:
            df = pd.read_excel(path, engine="openpyxl")
        # 컬럼명 문자열로 변환
        df.columns = [str(c) for c in df.columns]
        return df
    except Exception as e:
        raise ValueError(f"Excel 파일 파싱 실패: {e}")


def parse_word(path: str) -> pd.DataFrame:
    """Word 파일에서 첫 번째 테이블 추출"""
    try:
        from docx import Document
        doc = Document(path)
        if not doc.tables:
            raise ValueError("Word 파일에 테이블이 없습니다.")
        table = doc.tables[0]
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        if not rows:
            raise ValueError("Word 테이블이 비어 있습니다.")
        headers = rows[0]
        data = rows[1:]
        df = pd.DataFrame(data, columns=headers)
        return df
    except ImportError:
        raise ValueError("python-docx가 설치되어 있지 않습니다.")
    except Exception as e:
        raise ValueError(f"Word 파일 파싱 실패: {e}")


def parse_pptx(path: str) -> pd.DataFrame:
    """PPT 파일에서 첫 번째 테이블 shape 추출"""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        rows.append([cell.text.strip() for cell in row.cells])
                    if rows:
                        headers = rows[0]
                        data = rows[1:]
                        df = pd.DataFrame(data, columns=headers)
                        return df
        raise ValueError("PPT 파일에 테이블이 없습니다.")
    except ImportError:
        raise ValueError("python-pptx가 설치되어 있지 않습니다.")
    except Exception as e:
        raise ValueError(f"PPT 파일 파싱 실패: {e}")


def parse_file(path: str) -> pd.DataFrame:
    """파일 경로에 따라 적절한 파서 호출"""
    if not os.path.exists(path):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {path}")
    ext = Path(path).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"지원하지 않는 파일 형식입니다: {ext}")
    if ext in (".xlsx", ".xls"):
        return parse_excel(path)
    elif ext == ".docx":
        return parse_word(path)
    elif ext == ".pptx":
        return parse_pptx(path)
    else:
        raise ValueError(f"지원하지 않는 파일 형식입니다: {ext}")


def get_file_schema(path: str) -> Dict[str, Any]:
    """
    파일의 컬럼 목록과 샘플 데이터(최대 5행) 반환
    Returns: {"columns": [...], "sample": [[...], ...]}
    """
    df = parse_file(path)
    columns = list(df.columns)
    sample_df = df.head(5).fillna("")
    sample = []
    for _, row in sample_df.iterrows():
        sample.append([str(v) if v != "" else "" for v in row.tolist()])
    return {"columns": columns, "sample": sample}
