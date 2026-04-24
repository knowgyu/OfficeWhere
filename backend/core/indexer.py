import os
import re
import threading
import logging
from datetime import datetime, time as dt_time
from pathlib import Path
from typing import Any, Dict, List, Tuple

from ..database import (
    get_all_files,
    save_file_chunks,
    update_file_mtime,
    search_chunks,
    get_setting,
    set_setting,
)

logger = logging.getLogger(__name__)

_scheduler_thread: threading.Thread | None = None

_MAX_WORKERS = 8


# --- combined inspect + chunk (single file open) ---

def _inspect_and_chunk_excel(path: str) -> Tuple[List[str], List[Dict[str, str]]]:
    import pandas as pd

    chunks: List[Dict[str, str]] = []
    first_columns: List[str] = []

    try:
        xl = pd.ExcelFile(path, engine="openpyxl")
    except Exception:
        xl = pd.ExcelFile(path)

    for sheet_name in xl.sheet_names:
        df = xl.parse(sheet_name, dtype=str, keep_default_na=False)
        df.columns = [str(c) for c in df.columns]

        if not first_columns:
            first_columns = list(df.columns)

        prefix = f"{sheet_name} | " if len(xl.sheet_names) > 1 else ""
        header_text = " ".join(c for c in df.columns if c.strip())
        if header_text.strip():
            chunks.append({"location": f"{prefix}컬럼 헤더", "content": header_text})

        for row_idx, row in enumerate(df.itertuples(index=False), start=2):
            for col_name, cell_val in zip(df.columns, row):
                text = str(cell_val).strip()
                if text:
                    chunks.append({"location": f"{prefix}행 {row_idx}, 열 {col_name}", "content": text})

    return first_columns, chunks


def _inspect_and_chunk_word(path: str) -> Tuple[List[str], List[Dict[str, str]]]:
    from docx import Document

    doc = Document(path)
    chunks: List[Dict[str, str]] = []
    columns: List[str] = []
    para_idx = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            para_idx += 1
            chunks.append({"location": f"단락 {para_idx}", "content": text})

    for tbl_idx, table in enumerate(doc.tables, start=1):
        for r_idx, row in enumerate(table.rows, start=1):
            row_cells = [cell.text.strip() for cell in row.cells]
            if r_idx == 1 and tbl_idx == 1:
                columns = [c for c in row_cells if c]
            for c_idx, text in enumerate(row_cells, start=1):
                if text:
                    chunks.append({"location": f"표 {tbl_idx}, 행 {r_idx}, 열 {c_idx}", "content": text})

    return columns, chunks


def _inspect_and_chunk_pptx(path: str) -> Tuple[List[str], List[Dict[str, str]]]:
    from pptx import Presentation

    prs = Presentation(path)
    chunks: List[Dict[str, str]] = []
    columns: List[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    chunks.append({"location": f"슬라이드 {slide_num}", "content": text})
            if shape.has_table:
                for r_idx, row in enumerate(shape.table.rows, start=1):
                    for c_idx, cell in enumerate(row.cells, start=1):
                        text = cell.text_frame.text.strip()
                        if text:
                            if r_idx == 1 and not columns:
                                columns.append(text)
                            chunks.append({
                                "location": f"슬라이드 {slide_num}, 표 행 {r_idx} 열 {c_idx}",
                                "content": text,
                            })

    return columns, chunks


def inspect_and_chunk(path: str) -> Tuple[Dict[str, Any], List[Dict[str, str]]]:
    """파일 1회 파싱으로 컬럼 목록 + FTS 청크를 동시에 반환."""
    from .parser import get_file_type

    ext = Path(path).suffix.lower()
    if ext in (".xlsx", ".xls"):
        columns, chunks = _inspect_and_chunk_excel(path)
    elif ext == ".docx":
        columns, chunks = _inspect_and_chunk_word(path)
    elif ext == ".pptx":
        columns, chunks = _inspect_and_chunk_pptx(path)
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {ext}")

    return {
        "name": Path(path).name,
        "file_type": get_file_type(path),
        "columns": columns,
    }, chunks


# --- public API ---

def index_file(file_id: int, path: str) -> int:
    """파일을 파싱해 청크 인덱싱. 저장된 청크 수 반환."""
    ext = Path(path).suffix.lower()
    if ext in (".xlsx", ".xls"):
        _, chunks = _inspect_and_chunk_excel(path)
    elif ext == ".docx":
        _, chunks = _inspect_and_chunk_word(path)
    elif ext == ".pptx":
        _, chunks = _inspect_and_chunk_pptx(path)
    else:
        return 0

    save_file_chunks(file_id, chunks)
    update_file_mtime(file_id, os.path.getmtime(path))
    return len(chunks)


def _sanitize_fts_query(raw: str) -> str:
    """FTS5 MATCH 쿼리로 변환: 각 단어를 phrase로 AND 결합."""
    raw = re.sub(r'["\*\(\)\[\]\{\}\^~\?\\]', " ", raw)
    words = raw.split()
    if not words:
        return '""'
    return " ".join(f'"{w}"' for w in words)


def search(query: str, limit: int = 100) -> list:
    if not query.strip():
        return []
    fts_query = _sanitize_fts_query(query)
    return search_chunks(fts_query, limit=limit)


def reindex_all() -> Dict[str, int]:
    """등록된 모든 파일을 병렬 강제 재인덱싱."""
    from concurrent.futures import ThreadPoolExecutor

    files = get_all_files()

    def _reindex_one(f: Dict) -> str:
        path = f["path"]
        if not os.path.exists(path):
            return "failed"
        try:
            index_file(f["id"], path)
            return "success"
        except Exception as e:
            logger.warning("index_file failed for %s: %s", path, e)
            return "failed"

    with ThreadPoolExecutor(max_workers=_MAX_WORKERS) as executor:
        outcomes = list(executor.map(_reindex_one, files))

    set_setting("last_reindex_at", datetime.now().isoformat())
    return {
        "success": outcomes.count("success"),
        "failed": outcomes.count("failed"),
        "skipped": 0,
    }


def _do_reindex_incremental():
    """변경된 파일만 재인덱싱 (스케줄러 내부 호출용)."""
    files = get_all_files()
    for f in files:
        path = f["path"]
        if not os.path.exists(path):
            continue
        try:
            current_mtime = os.path.getmtime(path)
            stored_mtime = f.get("file_mtime")
            if stored_mtime is not None and abs(current_mtime - stored_mtime) < 1.0:
                continue
            index_file(f["id"], path)
        except Exception as e:
            logger.warning("incremental index failed for %s: %s", path, e)

    set_setting("last_reindex_at", datetime.now().isoformat())


def _scheduler_loop():
    import time

    while True:
        time.sleep(60)
        try:
            mode = get_setting("reindex_mode", "manual")
            if mode == "manual":
                continue

            if mode == "interval":
                interval_hours = float(get_setting("reindex_interval_hours", "24"))
                last_str = get_setting("last_reindex_at", "")
                if last_str:
                    last_dt = datetime.fromisoformat(last_str)
                    elapsed_hours = (datetime.now() - last_dt).total_seconds() / 3600
                    if elapsed_hours < interval_hours:
                        continue
                _do_reindex_incremental()

            elif mode == "daily":
                target_str = get_setting("reindex_daily_time", "03:00")
                now = datetime.now()
                try:
                    target_h, target_m = map(int, target_str.split(":"))
                except ValueError:
                    continue
                last_str = get_setting("last_reindex_at", "")
                # 오늘 목표 시각 이후이고 오늘 아직 실행 안 했으면 실행
                target_today = now.replace(hour=target_h, minute=target_m, second=0, microsecond=0)
                if now >= target_today:
                    if last_str:
                        last_dt = datetime.fromisoformat(last_str)
                        if last_dt.date() >= now.date():
                            continue
                    _do_reindex_incremental()
        except Exception as e:
            logger.warning("scheduler error: %s", e)


def start_scheduler():
    global _scheduler_thread
    if _scheduler_thread and _scheduler_thread.is_alive():
        return
    _scheduler_thread = threading.Thread(target=_scheduler_loop, daemon=True, name="reindex-scheduler")
    _scheduler_thread.start()
