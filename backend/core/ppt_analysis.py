from __future__ import annotations

from difflib import SequenceMatcher
import posixpath
import re
import zipfile
from typing import Any, Dict, List
from xml.etree import ElementTree as ET

from .normalizer import normalize_value

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
P_TAG = f"{{{P_NS}}}"
A_TAG = f"{{{A_NS}}}"
R_ID = f"{{{R_NS}}}id"
_SLIDE_NUMBER_RE = re.compile(r"slide(\d+)\.xml$")


def _safe_int(value: Any) -> int:
    if value is None:
        return 0
    try:
        return int(value)
    except (TypeError, ValueError):
        try:
            return int(float(str(value)))
        except (TypeError, ValueError):
            return 0


def _coerce_position(value: Any) -> int:
    if value is None:
        return 0
    try:
        return int(value)
    except (TypeError, ValueError):
        try:
            return int(float(str(value)))
        except (TypeError, ValueError):
            return 0


def _relationship_targets(archive: zipfile.ZipFile, rels_path: str, base_dir: str) -> Dict[str, str]:
    try:
        rels_xml = archive.read(rels_path)
    except KeyError:
        return {}
    root = ET.fromstring(rels_xml)
    targets: Dict[str, str] = {}
    for rel in root.findall(f"{{{PKG_REL_NS}}}Relationship"):
        rel_id = rel.get("Id")
        target = rel.get("Target")
        if not rel_id or not target:
            continue
        if target.startswith("/"):
            normalized = target.lstrip("/")
        else:
            normalized = posixpath.normpath(posixpath.join(base_dir, target))
        targets[rel_id] = normalized
    return targets


def _slide_paths(archive: zipfile.ZipFile) -> List[str]:
    rel_targets = _relationship_targets(
        archive,
        "ppt/_rels/presentation.xml.rels",
        "ppt",
    )
    paths: List[str] = []
    try:
        root = ET.fromstring(archive.read("ppt/presentation.xml"))
        for slide_id in root.findall(f".//{P_TAG}sldId"):
            rel_id = slide_id.get(R_ID)
            target = rel_targets.get(rel_id or "")
            if target and target in archive.namelist():
                paths.append(target)
    except (KeyError, ET.ParseError):
        paths = []

    if paths:
        return paths

    return sorted(
        (name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")),
        key=lambda name: int(_SLIDE_NUMBER_RE.search(name).group(1)) if _SLIDE_NUMBER_RE.search(name) else name,
    )


def _texts_in(element: ET.Element) -> List[str]:
    return [node.text for node in element.iter(f"{A_TAG}t") if node.text]


def _shape_text(shape: ET.Element) -> str:
    return "\n".join(text.strip() for text in _texts_in(shape) if text and text.strip()).strip()


def _shape_position(shape: ET.Element) -> tuple[int, int]:
    offset = shape.find(f".//{A_TAG}off")
    return (
        _safe_int(offset.get("y")) if offset is not None else 0,
        _safe_int(offset.get("x")) if offset is not None else 0,
    )


def _is_title_shape(shape: ET.Element) -> bool:
    for placeholder in shape.findall(f".//{P_TAG}ph"):
        placeholder_type = placeholder.get("type")
        if placeholder_type in {"title", "ctrTitle", "subTitle"}:
            return True
    non_visual_props = shape.find(f".//{P_TAG}cNvPr")
    name = (non_visual_props.get("name", "") if non_visual_props is not None else "").lower()
    return "title" in name or "제목" in name


def _table_items(shape: ET.Element, shape_idx: int) -> List[Dict[str, Any]]:
    items: List[Dict[str, Any]] = []
    for row_idx, row in enumerate(shape.findall(f".//{A_TAG}tr"), start=1):
        cell_texts: List[str] = []
        for cell in row.findall(f"{A_TAG}tc"):
            cell_text = "\n".join(text.strip() for text in _texts_in(cell) if text.strip()).strip()
            if cell_text:
                cell_texts.append(cell_text)
        row_text = " | ".join(cell_texts)
        if not row_text:
            continue
        items.append(
            {
                "item_type": "table_row",
                "location": f"shape:{shape_idx}/row:{row_idx}",
                "text": row_text,
                "normalized_text": normalize_value(row_text),
            }
        )
    return items


def _slide_elements(root: ET.Element) -> List[ET.Element]:
    shape_tree = root.find(f".//{P_TAG}spTree")
    if shape_tree is None:
        return []
    return [
        child
        for child in list(shape_tree)
        if child.tag in {f"{P_TAG}sp", f"{P_TAG}graphicFrame"}
    ]


def extract_ppt_slides(path: str) -> List[Dict[str, Any]]:
    """Extract slide text without reading embedded media parts.

    PPTX is a ZIP package.  We intentionally read only presentation and slide
    XML so broken/heavy videos (for example `ppt/media/*.mkv`) do not slow down
    or fail text indexing/version comparison.
    """
    slides: List[Dict[str, Any]] = []

    with zipfile.ZipFile(path) as archive:
        slide_paths = _slide_paths(archive)
        for slide_number, slide_path in enumerate(slide_paths, start=1):
            try:
                root = ET.fromstring(archive.read(slide_path))
            except (KeyError, ET.ParseError):
                continue

            sortable_items: List[tuple[int, int, Dict[str, Any]]] = []
            title = ""
            shapes = _slide_elements(root)
            for shape_idx, shape in enumerate(shapes, start=1):
                top, left = _shape_position(shape)
                if shape.find(f".//{A_TAG}tbl") is not None:
                    for item in _table_items(shape, shape_idx):
                        sortable_items.append((top, left, item))
                    if shape.tag == f"{P_TAG}graphicFrame":
                        continue

                text = _shape_text(shape)
                if not text:
                    continue
                if not title and _is_title_shape(shape):
                    title = text
                sortable_items.append(
                    (
                        top,
                        left,
                        {
                            "item_type": "text",
                            "location": f"shape:{shape_idx}",
                            "text": text,
                            "normalized_text": normalize_value(text),
                        },
                    )
                )

            sortable_items.sort(key=lambda item: (item[0], item[1], item[2]["location"]))
            items = [item[2] for item in sortable_items]
            if not title and items:
                title = items[0]["text"]
            signature_source = " ".join([title] + [item["text"] for item in items])
            slides.append(
                {
                    "slide_number": slide_number,
                    "title": title,
                    "signature": normalize_value(signature_source).lower(),
                    "items": items,
                }
            )

    return slides


def _extract_ppt_slides_with_package(path: str) -> List[Dict[str, Any]]:
    """Legacy package-based extractor kept for diagnostics; not used by default."""
    from pptx import Presentation

    presentation = Presentation(path)
    slides: List[Dict[str, Any]] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        sortable_items: List[tuple[int, int, Dict[str, Any]]] = []
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            top = _coerce_position(getattr(shape, "top", 0))
            left = _coerce_position(getattr(shape, "left", 0))
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    sortable_items.append(
                        (
                            top,
                            left,
                            {
                                "item_type": "text",
                                "location": f"shape:{shape_idx}",
                                "text": text,
                                "normalized_text": normalize_value(text),
                            },
                        )
                    )
            if getattr(shape, "has_table", False):
                for row_idx, row in enumerate(shape.table.rows, start=1):
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if not row_text:
                        continue
                    sortable_items.append(
                        (
                            top,
                            left,
                            {
                                "item_type": "table_row",
                                "location": f"shape:{shape_idx}/row:{row_idx}",
                                "text": row_text,
                                "normalized_text": normalize_value(row_text),
                            },
                        )
                    )

        sortable_items.sort(key=lambda item: (item[0], item[1], item[2]["location"]))
        items = [item[2] for item in sortable_items]
        title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else ""
        if not title and items:
            title = items[0]["text"]
        signature_source = " ".join([title] + [item["text"] for item in items])
        slides.append(
            {
                "slide_number": slide_number,
                "title": title,
                "signature": normalize_value(signature_source).lower(),
                "items": items,
            }
        )

    return slides


def inspect_ppt_slides(slides: List[Dict[str, Any]]) -> Dict[str, Any]:
    sample = [[slide["slide_number"], slide["title"], len(slide["items"])] for slide in slides[:5]]
    return {
        "parser_config": {},
        "table_candidates": [],
        "columns": ["slide_number", "title", "item_count"],
        "sample": sample,
        "slide_count": len(slides),
    }


def slide_similarity(left_slide: Dict[str, Any], right_slide: Dict[str, Any]) -> float:
    if left_slide["title"] and left_slide["title"] == right_slide["title"]:
        return 1.0
    return SequenceMatcher(None, left_slide["signature"], right_slide["signature"]).ratio()


def inspect_ppt_file(path: str) -> Dict[str, Any]:
    return inspect_ppt_slides(extract_ppt_slides(path))
