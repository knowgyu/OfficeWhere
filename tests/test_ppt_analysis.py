import zipfile

from pptx import Presentation
from pptx.util import Inches

from backend.core.ppt_analysis import _coerce_position, extract_ppt_slides


def test_coerce_position_accepts_float_like_string():
    assert _coerce_position("3520440.0") == 3520440


def test_coerce_position_falls_back_for_invalid_value():
    assert _coerce_position("not-a-number") == 0


def test_extract_ppt_slides_does_not_read_embedded_media(monkeypatch, tmp_path):
    path = tmp_path / "with-media.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Overview"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4), Inches(1))
    textbox.text_frame.text = "본문"
    presentation.save(path)

    with zipfile.ZipFile(path, "a") as archive:
        archive.writestr("ppt/media/media8.mkv", b"placeholder")

    real_read = zipfile.ZipFile.read

    def guarded_read(self, name, *args, **kwargs):
        assert not str(name).startswith("ppt/media/")
        return real_read(self, name, *args, **kwargs)

    monkeypatch.setattr(zipfile.ZipFile, "read", guarded_read)

    slides = extract_ppt_slides(str(path))

    assert slides[0]["title"] == "Overview"
    assert any(item["text"] == "본문" for item in slides[0]["items"])


def test_extract_ppt_slides_includes_graphic_frame_table_rows(tmp_path):
    path = tmp_path / "table.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Table"
    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1.5), Inches(4), Inches(1.2))
    table = table_shape.table
    table.cell(0, 0).text = "항목"
    table.cell(0, 1).text = "값"
    table.cell(1, 0).text = "예산"
    table.cell(1, 1).text = "100"
    presentation.save(path)

    slides = extract_ppt_slides(str(path))

    table_rows = [item["text"] for item in slides[0]["items"] if item["item_type"] == "table_row"]
    assert "항목 | 값" in table_rows
    assert "예산 | 100" in table_rows
