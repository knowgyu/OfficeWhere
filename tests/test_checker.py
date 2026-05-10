import os
from pathlib import Path

import pytest
from docx import Document
from fastapi import HTTPException
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from backend.api.check import consistency_check
from backend.core.checker import run_consistency_check
from backend.core.excel_diff_grid import build_excel_diff_grid
from backend.core.file_access import inspect_file_path
from backend.models.schemas import CheckRequest


def _write_excel_with_offset_table(path: Path):
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Sheet1"
    worksheet.title = "사업현황"
    worksheet["A1"] = "2026 사업 목록"
    worksheet["C3"] = "과제명"
    worksheet["D3"] = "담당자"
    worksheet["E3"] = "예산"
    worksheet["C4"] = "A"
    worksheet["D4"] = "Kim"
    worksheet["E4"] = "100"
    worksheet["C5"] = "B"
    worksheet["D5"] = "Lee"
    worksheet["E5"] = "200"
    workbook.save(path)


def _write_excel_with_offset_conflict_table(path: Path, budget: str):
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "사업현황"
    worksheet["A1"] = "2026 사업 목록"
    worksheet["C3"] = "과제명"
    worksheet["D3"] = "담당자"
    worksheet["E3"] = "예산"
    # Row 4 is intentionally blank. The parser keeps source coordinates, so
    # location metadata must preserve the blank-row offset.
    worksheet["C5"] = "A"
    worksheet["D5"] = "Kim"
    worksheet["E5"] = budget
    worksheet["C6"] = "B"
    worksheet["D6"] = "Lee"
    worksheet["E6"] = "200"
    workbook.save(path)


def _write_tabular_excel(path: Path, data: dict):
    workbook = Workbook()
    worksheet = workbook.active
    headers = list(data.keys())
    worksheet.append(headers)
    row_count = max((len(values) for values in data.values()), default=0)
    for row_index in range(row_count):
        worksheet.append([
            values[row_index] if row_index < len(values) else ""
            for values in data.values()
        ])
    workbook.save(path)


def _write_multisheet_excel(path: Path, detail_value: str):
    workbook = Workbook()
    workbook.active.title = "요약"
    workbook.active["A1"] = "공통요약"
    detail = workbook.create_sheet("세부")
    detail["B2"] = detail_value
    workbook.save(path)


def _write_multisheet_grid_excel(path: Path, summary_value: str, detail_value: str):
    workbook = Workbook()
    summary = workbook.active
    summary.title = "요약"
    summary["A1"] = summary_value
    detail = workbook.create_sheet("세부")
    detail["B2"] = detail_value
    workbook.save(path)


def _write_word(path: Path, paragraph_text: str, table_value: str):
    document = Document()
    document.add_paragraph("공통 서론")
    document.add_paragraph(paragraph_text)
    table = document.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "ID"
    table.rows[0].cells[1].text = "Value"
    table.rows[1].cells[0].text = "1"
    table.rows[1].cells[1].text = table_value
    document.save(path)


def _write_word_with_second_page_change(path: Path, second_page_text: str):
    document = Document()
    document.add_paragraph("첫 페이지 공통")
    document.add_page_break()
    document.add_paragraph(second_page_text)
    document.save(path)


def _add_textbox(slide, left: float, top: float, text: str):
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(4), Inches(1))
    textbox.text_frame.text = text


def _write_ppt(path: Path, first_body: str, include_inserted_slide: bool):
    presentation = Presentation()
    layout = presentation.slide_layouts[5]

    slide1 = presentation.slides.add_slide(layout)
    slide1.shapes.title.text = "Overview"
    _add_textbox(slide1, 1, 1.5, first_body)

    if include_inserted_slide:
        inserted = presentation.slides.add_slide(layout)
        inserted.shapes.title.text = "Inserted"
        _add_textbox(inserted, 1, 1.5, "추가된 슬라이드")

    slide2 = presentation.slides.add_slide(layout)
    slide2.shapes.title.text = "Plan"
    _add_textbox(slide2, 1, 1.5, "기존 슬라이드")

    presentation.save(path)


def test_excel_inspect_uses_visible_range_without_table_detection(tmp_path):
    file_path = tmp_path / "offset.xlsx"
    _write_excel_with_offset_table(file_path)

    result = inspect_file_path(str(file_path))

    assert result["columns"][0] == "2026 사업 목록"
    flattened_sample = [cell for row in result["sample"] for cell in row]
    assert {"과제명", "담당자", "예산", "A", "Kim", "100"} <= set(flattened_sample)


def test_excel_consistency_uses_cell_diffs_not_registered_keys(tmp_path):
    file_a = tmp_path / "a.xlsx"
    file_b = tmp_path / "b.xlsx"
    _write_tabular_excel(file_a, {"과제명": ["A", "B"], "예산": ["100", "200"], "담당자": ["Kim", "Lee"]})
    _write_tabular_excel(file_b, {"과제명": ["A", "C"], "예산": ["999", "300"]})

    result = run_consistency_check(
        [
            {'id': 1, 'path': str(file_a), 'name': 'a.xlsx', 'file_type': 'Excel'},
            {'id': 2, 'path': str(file_b), 'name': 'b.xlsx', 'file_type': 'Excel'},
        ]
    )

    assert result["mode"] == "excel"
    issue_types = {issue["issue_type"] for issue in result["excel"]["issues"]}
    assert {"value_conflict", "value_removed"} <= issue_types
    assert "missing_column" not in issue_types
    assert "missing_key" not in issue_types

    conflict = next(issue for issue in result["excel"]["issues"] if issue["issue_type"] == "value_conflict")
    assert conflict["key"] == "2"
    assert conflict["column"] == "B"
    values_by_file = {entry["file_id"]: entry for entry in conflict["values"]}
    assert values_by_file[1]["row_numbers"] == [2]
    assert values_by_file[1]["column_letters"] == ["B"]
    assert values_by_file[1]["cell_refs"] == ["B2"]
    assert values_by_file[1]["row_count"] == 1
    assert values_by_file[2]["row_numbers"] == [2]
    assert values_by_file[2]["column_letters"] == ["B"]
    assert values_by_file[2]["cell_refs"] == ["B2"]
    assert values_by_file[2]["row_count"] == 1


def test_excel_consistency_reports_cell_value_added_and_removed(tmp_path):
    file_a = tmp_path / "cell-a.xlsx"
    file_b = tmp_path / "cell-b.xlsx"
    _write_tabular_excel(
        file_a,
        {"과제명": ["A", "B"], "담당자": ["", "Kim"], "예산": ["100", "200"]},
    )
    _write_tabular_excel(
        file_b,
        {"과제명": ["A", "B"], "담당자": ["Lee", ""], "예산": ["100", "200"]},
    )

    result = run_consistency_check(
        [
            {'id': 1, 'path': str(file_a), 'name': 'cell-a.xlsx', 'file_type': 'Excel'},
            {'id': 2, 'path': str(file_b), 'name': 'cell-b.xlsx', 'file_type': 'Excel'},
        ]
    )

    issues = result["excel"]["issues"]
    added = next(issue for issue in issues if issue["issue_type"] == "value_added")
    removed = next(issue for issue in issues if issue["issue_type"] == "value_removed")

    assert added["key"] == "2"
    assert added["column"] == "B"
    assert added["message"] == "2행 B열 값이 추가되었습니다."
    assert [entry["values"] for entry in added["values"]] == [["(빈 값)"], ["Lee"]]
    assert added["values"][1]["cell_refs"] == ["B2"]

    assert removed["key"] == "3"
    assert removed["column"] == "B"
    assert removed["message"] == "3행 B열 값이 삭제되었습니다."
    assert [entry["values"] for entry in removed["values"]] == [["Kim"], ["(빈 값)"]]
    assert removed["values"][0]["cell_refs"] == ["B3"]


def test_excel_consistency_compares_all_visible_sheets(tmp_path):
    file_a = tmp_path / "multi-a.xlsx"
    file_b = tmp_path / "multi-b.xlsx"
    _write_multisheet_excel(file_a, "기존상세")
    _write_multisheet_excel(file_b, "변경상세")

    result = run_consistency_check(
        [
            {'id': 1, 'path': str(file_a), 'name': 'multi-a.xlsx', 'file_type': 'Excel'},
            {'id': 2, 'path': str(file_b), 'name': 'multi-b.xlsx', 'file_type': 'Excel'},
        ]
    )

    issue = result["excel"]["issues"][0]
    assert issue["sheet_name"] == "세부"
    assert issue["key"] == "세부!2"
    assert issue["column"] == "B"
    assert issue["message"] == "세부 시트 | 2행 B열 값이 변경되었습니다."
    assert issue["values"][0]["cell_refs"] == ["세부!B2"]
    assert issue["values"][0]["sheet_name"] == "세부"


def test_excel_consistency_uses_fresh_indexed_sheet_cache(tmp_path, monkeypatch):
    from backend.core.indexer import inspect_and_chunk
    from backend.database import get_file_by_id, init_db, save_indexed_file

    monkeypatch.setattr("backend.database.DB_PATH", tmp_path / "test.db")
    monkeypatch.setattr("backend.database.DB_DIR", tmp_path)
    init_db()

    file_a = tmp_path / "indexed-a.xlsx"
    file_b = tmp_path / "indexed-b.xlsx"
    _write_multisheet_excel(file_a, "캐시전")
    _write_multisheet_excel(file_b, "캐시후")

    file_ids = []
    for file_path in (file_a, file_b):
        info, chunks = inspect_and_chunk(str(file_path))
        file_id = save_indexed_file(path=str(file_path), name=info['name'], file_type=info['file_type'], column_count=len(info['columns']), chunks=chunks, file_mtime=file_path.stat().st_mtime, excel_sheets=info['excel_sheets'], excel_cells=info['excel_cells'])
        file_ids.append(file_id)

    def fail_source_parse(_path):
        raise AssertionError("fresh indexed Excel comparison should not parse source sheets")

    monkeypatch.setattr("backend.core.excel_compare.extract_excel_used_ranges", fail_source_parse)
    file_infos = [
        {
            "id": file_id,
            "path": get_file_by_id(file_id)["path"],
            "name": get_file_by_id(file_id)["name"],
            "file_type": get_file_by_id(file_id)["file_type"],
            "file_mtime": get_file_by_id(file_id)["file_mtime"],
        }
        for file_id in file_ids
    ]

    result = run_consistency_check(file_infos)

    issue = result["excel"]["issues"][0]
    assert issue["sheet_name"] == "세부"
    assert [entry["values"][0] for entry in issue["values"]] == ["캐시전", "캐시후"]
    assert result["metadata"]["warnings"] == []


def test_excel_indexed_payload_used_when_source_mtime_newer_with_warning(tmp_path, monkeypatch):
    from backend.core.indexer import inspect_and_chunk
    from backend.database import get_file_by_id, init_db, save_indexed_file

    monkeypatch.setattr("backend.database.DB_PATH", tmp_path / "test.db")
    monkeypatch.setattr("backend.database.DB_DIR", tmp_path)
    init_db()

    file_a = tmp_path / "indexed-stale-a.xlsx"
    file_b = tmp_path / "indexed-stale-b.xlsx"
    _write_multisheet_excel(file_a, "색인전")
    _write_multisheet_excel(file_b, "색인후")

    file_ids = []
    for file_path in (file_a, file_b):
        info, chunks = inspect_and_chunk(str(file_path))
        file_id = save_indexed_file(path=str(file_path), name=info['name'], file_type=info['file_type'], column_count=len(info['columns']), chunks=chunks, file_mtime=file_path.stat().st_mtime, excel_sheets=info['excel_sheets'], excel_cells=info['excel_cells'])
        file_ids.append(file_id)

    newer_mtime = file_b.stat().st_mtime + 5
    os.utime(file_b, (newer_mtime, newer_mtime))

    def fail_source_parse(_path):
        raise AssertionError("newer source mtime should warn, not force source parse")

    monkeypatch.setattr("backend.core.excel_compare.extract_excel_used_ranges", fail_source_parse)
    file_infos = [
        {
            "id": file_id,
            "path": get_file_by_id(file_id)["path"],
            "name": get_file_by_id(file_id)["name"],
            "file_type": get_file_by_id(file_id)["file_type"],
            "file_mtime": get_file_by_id(file_id)["file_mtime"],
        }
        for file_id in file_ids
    ]

    result = run_consistency_check(file_infos)

    issue = result["excel"]["issues"][0]
    assert [entry["values"][0] for entry in issue["values"]] == ["색인전", "색인후"]
    assert any(warning["type"] == "source_may_be_newer" for warning in result["metadata"]["warnings"])


def test_excel_sparse_diff_uses_coordinate_union_not_rectangle(monkeypatch):
    from backend.core import excel_compare

    lookup_count = 0

    class CountingCells(dict):
        def get(self, key, default=None):
            nonlocal lookup_count
            lookup_count += 1
            return super().get(key, default)

    file_infos = [
        {"id": 1, "path": "/tmp/before.xlsx", "name": "before.xlsx", "file_type": "Excel"},
        {"id": 2, "path": "/tmp/after.xlsx", "name": "after.xlsx", "file_type": "Excel"},
    ]
    payloads = [
        {
            "info": file_infos[0],
            "sheets": {"Sheet1": {"sheet_index": 1, "row_count": 10_000, "column_count": 26}},
            "cells": CountingCells({("Sheet1", 1, 1): "공통", ("Sheet1", 10_000, 26): "이전"}),
        },
        {
            "info": file_infos[1],
            "sheets": {"Sheet1": {"sheet_index": 1, "row_count": 10_000, "column_count": 26}},
            "cells": CountingCells({("Sheet1", 1, 1): "공통", ("Sheet1", 10_000, 26): "이후"}),
        },
    ]

    monkeypatch.setattr(excel_compare, "_excel_payloads", lambda _infos: (payloads, excel_compare._default_compare_metadata()))

    result = excel_compare.compare_excel_versions_by_cells(file_infos)

    assert lookup_count == 4
    assert result["total_keys"] == 2
    assert result["metadata"]["compared_cell_count"] == 2
    assert result["metadata"]["changed_cell_count"] == 1
    assert result["issues"][0]["key"] == "10000"
    assert result["issues"][0]["column"] == "Z"


def test_excel_sparse_metadata_warnings_are_not_fake_issues(monkeypatch):
    from backend.core import excel_compare

    file_infos = [
        {"id": 1, "path": "/tmp/before.xlsx", "name": "before.xlsx", "file_type": "Excel"},
        {"id": 2, "path": "/tmp/after.xlsx", "name": "after.xlsx", "file_type": "Excel"},
    ]
    before_cells = {("Sheet1", row, 1): f"이전-{row}" for row in range(1, 502)}
    after_cells = {("Sheet1", row, 1): f"이후-{row}" for row in range(1, 502)}
    payloads = [
        {
            "info": file_infos[0],
            "sheets": {"Sheet1": {"sheet_index": 1, "row_count": 501, "column_count": 1}},
            "cells": before_cells,
        },
        {
            "info": file_infos[1],
            "sheets": {"Sheet1": {"sheet_index": 1, "row_count": 501, "column_count": 1}},
            "cells": after_cells,
        },
    ]

    monkeypatch.setattr(excel_compare, "_excel_payloads", lambda _infos: (payloads, excel_compare._default_compare_metadata()))

    result = excel_compare.compare_excel_versions_by_cells(file_infos)

    warning_types = {warning["type"] for warning in result["metadata"]["warnings"]}
    assert {"truncated", "high_change_ratio"} <= warning_types
    assert len(result["issues"]) == excel_compare.EXCEL_VERSION_CELL_ISSUE_LIMIT
    assert all(issue["key"] != "truncated" for issue in result["issues"])
    assert result["metadata"]["compared_cell_count"] == 501
    assert result["metadata"]["changed_cell_count"] == 501


def test_excel_consistency_reports_offset_cell_refs_after_blank_rows(tmp_path):
    file_a = tmp_path / "offset-a.xlsx"
    file_b = tmp_path / "offset-b.xlsx"
    _write_excel_with_offset_conflict_table(file_a, "100")
    _write_excel_with_offset_conflict_table(file_b, "999")

    result = run_consistency_check(
        [
            {'id': 1, 'path': str(file_a), 'name': 'offset-a.xlsx', 'file_type': 'Excel'},
            {'id': 2, 'path': str(file_b), 'name': 'offset-b.xlsx', 'file_type': 'Excel'},
        ]
    )

    conflict = next(issue for issue in result["excel"]["issues"] if issue["issue_type"] == "value_conflict")
    assert conflict["key"] == "5"
    assert conflict["column"] == "E"
    for entry in conflict["values"]:
        assert entry["row_numbers"] == [5]
        assert entry["column_letters"] == ["E"]
        assert entry["cell_refs"] == ["E5"]
        assert entry["row_count"] == 1


def test_excel_version_history_uses_current_used_range(tmp_path):
    previous = tmp_path / "budget-v1.xlsx"
    latest = tmp_path / "budget-v2.xlsx"
    _write_tabular_excel(previous, {"과제명": ["A"], "예산": ["100"]})
    _write_tabular_excel(latest, {"과제명": ["A"], "예산": ["999"]})

    result = run_consistency_check([{'id': 1, 'path': str(previous), 'name': 'budget-v1.xlsx', 'file_type': 'Excel'}, {'id': 2, 'path': str(latest), 'name': 'budget-v2.xlsx', 'file_type': 'Excel'}])

    issue = next(issue for issue in result["excel"]["issues"] if issue["issue_type"] == "value_conflict")
    assert issue["key"] == "2"
    assert issue["column"] == "B"
    assert issue["message"] == "2행 B열 값이 변경되었습니다."
    assert [entry["cell_refs"] for entry in issue["values"]] == [["B2"], ["B2"]]
    assert [entry["values"] for entry in issue["values"]] == [["100"], ["999"]]


def test_excel_version_history_uses_row_coordinates(tmp_path):
    previous = tmp_path / "budget-v1.xlsx"
    latest = tmp_path / "budget-v2.xlsx"
    _write_tabular_excel(previous, {"과제명": ["A"], "예산": ["100"]})
    _write_tabular_excel(latest, {"과제명": ["A"], "예산": ["120"]})

    result = run_consistency_check([{'id': 1, 'path': str(previous), 'name': 'budget-v1.xlsx', 'file_type': 'Excel'}, {'id': 2, 'path': str(latest), 'name': 'budget-v2.xlsx', 'file_type': 'Excel'}])

    issue = next(issue for issue in result["excel"]["issues"] if issue["issue_type"] == "value_conflict")
    assert issue["key"] == "2"
    assert issue["column"] == "B"
    assert [entry["values"] for entry in issue["values"]] == [["100"], ["120"]]


def test_excel_diff_grid_returns_full_small_latest_sheet(tmp_path):
    latest = tmp_path / "grid-latest.xlsx"
    previous = tmp_path / "grid-previous.xlsx"
    _write_tabular_excel(latest, {"과제명": ["A", "B"], "담당자": ["Kim", "Lee"], "예산": ["100", "250"]})
    _write_tabular_excel(previous, {"과제명": ["A", "B"], "담당자": ["Kim", "Lee"], "예산": ["100", "200"]})

    result = build_excel_diff_grid(
        [
            {'id': 2, 'path': str(latest), 'name': 'grid-latest.xlsx', 'file_type': 'Excel'},
            {'id': 1, 'path': str(previous), 'name': 'grid-previous.xlsx', 'file_type': 'Excel'},
        ],
        [
            {
                "key": "3",
                "column": "C",
                "change_type": "changed",
                "histories": [
                    {
                        "change_type": "changed",
                        "from_file_id": 1,
                        "from_file_name": "grid-previous.xlsx",
                        "to_file_id": 2,
                        "to_file_name": "grid-latest.xlsx",
                        "before": "200",
                        "after": "250",
                        "label": "grid-previous.xlsx → grid-latest.xlsx",
                    }
                ],
            }
        ],
    )

    assert result["partial"] is False
    assert result["row_count"] == 5
    assert result["column_count"] == 5
    section = result["sections"][0]
    assert section["row_start"] == 1
    assert section["row_end"] == 5
    highlighted = [
        cell
        for row in section["rows"]
        for cell in row["cells"]
        if cell["highlight"] == "changed"
    ]
    assert len(highlighted) == 1
    assert highlighted[0]["column_name"] == "C"
    assert highlighted[0]["column_letter"] == "C"
    assert highlighted[0]["row_number"] == 3
    assert highlighted[0]["value"] == "250"
    assert highlighted[0]["histories"][0]["before"] == "200"


def test_excel_diff_grid_keeps_multi_sheet_sections_separable(tmp_path):
    latest = tmp_path / "multi-grid-latest.xlsx"
    previous = tmp_path / "multi-grid-previous.xlsx"
    _write_multisheet_grid_excel(latest, "요약 최신", "세부 최신")
    _write_multisheet_grid_excel(previous, "요약 이전", "세부 이전")

    result = build_excel_diff_grid(
        [
            {'id': 2, 'path': str(latest), 'name': 'multi-grid-latest.xlsx', 'file_type': 'Excel'},
            {'id': 1, 'path': str(previous), 'name': 'multi-grid-previous.xlsx', 'file_type': 'Excel'},
        ],
        [
            {
                "sheet_name": "요약",
                "key": "1",
                "column": "A",
                "change_type": "changed",
                "histories": [
                    {
                        "change_type": "changed",
                        "from_file_id": 1,
                        "from_file_name": "multi-grid-previous.xlsx",
                        "to_file_id": 2,
                        "to_file_name": "multi-grid-latest.xlsx",
                        "before": "요약 이전",
                        "after": "요약 최신",
                        "label": "multi-grid-previous.xlsx → multi-grid-latest.xlsx",
                    }
                ],
            },
            {
                "sheet_name": "세부",
                "key": "2",
                "column": "B",
                "change_type": "changed",
                "histories": [
                    {
                        "change_type": "changed",
                        "from_file_id": 1,
                        "from_file_name": "multi-grid-previous.xlsx",
                        "to_file_id": 2,
                        "to_file_name": "multi-grid-latest.xlsx",
                        "before": "세부 이전",
                        "after": "세부 최신",
                        "label": "multi-grid-previous.xlsx → multi-grid-latest.xlsx",
                    }
                ],
            },
        ],
    )

    assert result["sheet_name"] == "여러 시트"
    section_sheet_names = [section["sheet_name"] for section in result["sections"]]
    assert section_sheet_names == ["요약", "세부"]

    values_by_sheet = {
        section["sheet_name"]: [
            cell["value"]
            for row in section["rows"]
            for cell in row["cells"]
            if cell["highlight"] == "changed"
        ]
        for section in result["sections"]
    }
    assert values_by_sheet["요약"] == ["요약 최신"]
    assert values_by_sheet["세부"] == ["세부 최신"]


def test_excel_diff_grid_colors_only_latest_vs_previous_but_keeps_older_history(tmp_path):
    v1 = tmp_path / "budget-v1.xlsx"
    v2 = tmp_path / "budget-v2.xlsx"
    v3 = tmp_path / "budget-v3.xlsx"
    _write_tabular_excel(v1, {"ID": ["A", "B"], "값": ["초안", "유지"]})
    _write_tabular_excel(v2, {"ID": ["A", "B"], "값": ["중간", "유지"]})
    _write_tabular_excel(v3, {"ID": ["A", "B"], "값": ["중간", "최신"]})

    result = build_excel_diff_grid(
        [
            {'id': 3, 'path': str(v3), 'name': 'budget-v3.xlsx', 'file_type': 'Excel'},
            {'id': 2, 'path': str(v2), 'name': 'budget-v2.xlsx', 'file_type': 'Excel'},
            {'id': 1, 'path': str(v1), 'name': 'budget-v1.xlsx', 'file_type': 'Excel'},
        ],
        [],
    )

    cells = {
        (cell["row_number"], cell["column_letter"]): cell
        for row in result["sections"][0]["rows"]
        for cell in row["cells"]
    }

    older_change = cells[(2, "B")]
    latest_change = cells[(3, "B")]

    assert older_change["highlight"] is None
    assert older_change["histories"][0]["label"] == "budget-v1.xlsx → budget-v2.xlsx"
    assert older_change["histories"][0]["before"] == "초안"
    assert older_change["histories"][0]["after"] == "중간"

    assert latest_change["highlight"] == "changed"
    assert latest_change["histories"][0]["label"] == "budget-v2.xlsx → budget-v3.xlsx"
    assert latest_change["histories"][0]["before"] == "유지"
    assert latest_change["histories"][0]["after"] == "최신"


def test_excel_diff_grid_uses_largest_compared_range_for_removed_cells(tmp_path):
    latest = tmp_path / "smaller-latest.xlsx"
    previous = tmp_path / "larger-previous.xlsx"
    _write_tabular_excel(latest, {"ID": ["A", "B"], "담당자": ["Kim", "Lee"], "예산": ["100", "250"]})
    _write_tabular_excel(
        previous,
        {
            "ID": ["A", "B", "C", "D"],
            "담당자": ["Kim", "Lee", "Park", "Choi"],
            "예산": ["100", "200", "300", "400"],
            "상태": ["진행", "완료", "보류", "종료"],
            "비고": ["", "", "", "삭제될 값"],
        },
    )

    result = build_excel_diff_grid(
        [
            {'id': 2, 'path': str(latest), 'name': 'smaller-latest.xlsx', 'file_type': 'Excel'},
            {'id': 1, 'path': str(previous), 'name': 'larger-previous.xlsx', 'file_type': 'Excel'},
        ],
        [
            {
                "key": "5",
                "column": "E",
                "change_type": "removed",
                "histories": [
                    {
                        "change_type": "removed",
                        "from_file_id": 1,
                        "from_file_name": "larger-previous.xlsx",
                        "to_file_id": 2,
                        "to_file_name": "smaller-latest.xlsx",
                        "before": "삭제될 값",
                        "after": "",
                        "label": "larger-previous.xlsx → smaller-latest.xlsx",
                    }
                ],
            }
        ],
    )

    assert result["row_count"] == 7
    assert result["column_count"] == 7
    assert result["omitted_focus_count"] == 0
    highlighted = [
        cell
        for row in result["sections"][0]["rows"]
        for cell in row["cells"]
        if cell["highlight"] == "removed" and cell["row_number"] == 5 and cell["column_letter"] == "E"
    ]
    assert len(highlighted) == 1
    assert highlighted[0]["row_number"] == 5
    assert highlighted[0]["column_letter"] == "E"
    assert highlighted[0]["value"] == ""
    assert highlighted[0]["histories"][0]["before"] == "삭제될 값"


def test_excel_diff_grid_shows_small_compared_used_range_with_margin(tmp_path):
    latest = tmp_path / "dirty-latest.xlsx"
    previous = tmp_path / "dirty-previous.xlsx"
    for path, value in [(latest, "변경 후"), (previous, "변경 전")]:
        workbook = Workbook()
        worksheet = workbook.active
        worksheet.title = "Sheet1"
        worksheet["A1"] = "ID"
        worksheet["B1"] = "값"
        worksheet["A2"] = "A"
        worksheet["B2"] = value
        # Even when the focus list is short/truncated, 표로 보기 should use the
        # compared files' full used range plus a small visual margin.
        if path == previous:
            worksheet["AD50"] = "삭제될 먼 위치 값"
        workbook.save(path)

    result = build_excel_diff_grid(
        [
            {'id': 2, 'path': str(latest), 'name': 'dirty-latest.xlsx', 'file_type': 'Excel'},
            {'id': 1, 'path': str(previous), 'name': 'dirty-previous.xlsx', 'file_type': 'Excel'},
        ],
        [
            {
                "key": "2",
                "column": "B",
                "change_type": "changed",
                "histories": [
                    {
                        "change_type": "changed",
                        "from_file_id": 1,
                        "from_file_name": "dirty-previous.xlsx",
                        "to_file_id": 2,
                        "to_file_name": "dirty-latest.xlsx",
                        "before": "변경 전",
                        "after": "변경 후",
                        "label": "dirty-previous.xlsx → dirty-latest.xlsx",
                    }
                ],
            }
        ],
    )

    assert result["row_count"] == 52
    assert result["column_count"] == 32
    assert result["partial"] is False
    section = result["sections"][0]
    assert section["row_start"] == 1
    assert section["row_end"] == 52
    assert section["col_start"] == 1
    assert section["col_end"] == 32
    highlighted = [
        cell
        for row in section["rows"]
        for cell in row["cells"]
        if cell["highlight"] == "changed"
    ]
    assert len(highlighted) == 1
    assert highlighted[0]["row_number"] == 2
    assert highlighted[0]["column_letter"] == "B"
    assert highlighted[0]["value"] == "변경 후"
    removed = [
        cell
        for row in section["rows"]
        for cell in row["cells"]
        if cell["highlight"] == "removed" and cell["row_number"] == 50 and cell["column_letter"] == "AD"
    ]
    assert len(removed) == 1
    assert removed[0]["value"] == ""
    assert removed[0]["histories"][0]["before"] == "삭제될 먼 위치 값"


def test_excel_diff_grid_limits_large_far_focus_and_keeps_context_column(tmp_path):
    latest = tmp_path / "large-latest.xlsx"
    previous = tmp_path / "large-previous.xlsx"
    rows = 120
    data = {"ID": [f"K{row}" for row in range(1, rows + 1)]}
    for column in range(1, 120):
        data[f"C{column}"] = [f"{row}-{column}" for row in range(1, rows + 1)]
    _write_tabular_excel(latest, data)
    _write_tabular_excel(previous, data)

    result = build_excel_diff_grid(
        [
            {'id': 2, 'path': str(latest), 'name': 'large-latest.xlsx', 'file_type': 'Excel'},
            {'id': 1, 'path': str(previous), 'name': 'large-previous.xlsx', 'file_type': 'Excel'},
        ],
        [{"key": "101", "column": "CW", "change_type": "added", "histories": []}],
    )

    assert result["partial"] is True
    section = result["sections"][0]
    assert section["row_start"] > 1
    assert section["row_end"] - section["row_start"] + 1 < rows + 1
    assert section["col_start"] > 1
    assert any(column["name"] == "A" for column in section["columns"])
    highlighted = [
        cell
        for row in section["rows"]
        for cell in row["cells"]
        if cell["highlight"] == "added"
    ]
    assert highlighted
    assert highlighted[0]["column_name"] == "CW"


def test_excel_diff_grid_resolves_partial_focus_by_row_number_and_header(tmp_path):
    latest = tmp_path / "large-key-latest.xlsx"
    previous = tmp_path / "large-key-previous.xlsx"
    rows = 120
    data = {"ID": [f"K{row}" for row in range(1, rows + 1)]}
    for column in range(1, 120):
        data[f"C{column}"] = [f"{row}-{column}" for row in range(1, rows + 1)]

    previous_data = {key: list(values) for key, values in data.items()}
    latest_data = {key: list(values) for key, values in data.items()}
    previous_data["C119"][100] = "old-budget"
    latest_data["C119"][100] = "new-budget"
    _write_tabular_excel(latest, latest_data)
    _write_tabular_excel(previous, previous_data)

    result = build_excel_diff_grid(
        [
            {'id': 2, 'path': str(latest), 'name': 'large-key-latest.xlsx', 'file_type': 'Excel'},
            {'id': 1, 'path': str(previous), 'name': 'large-key-previous.xlsx', 'file_type': 'Excel'},
        ],
        [
            {
                "key": "102",
                "column": "C119",
                "change_type": "changed",
                "histories": [
                    {
                        "change_type": "changed",
                        "from_file_id": 1,
                        "from_file_name": "large-key-previous.xlsx",
                        "to_file_id": 2,
                        "to_file_name": "large-key-latest.xlsx",
                        "before": "old-budget",
                        "after": "new-budget",
                        "label": "large-key-previous.xlsx → large-key-latest.xlsx",
                    }
                ],
            }
        ],
    )

    assert result["partial"] is True
    assert result["omitted_focus_count"] == 0
    highlighted = [
        cell
        for section in result["sections"]
        for row in section["rows"]
        for cell in row["cells"]
        if cell["highlight"] == "changed"
    ]
    assert len(highlighted) == 1
    assert highlighted[0]["row_number"] == 102
    assert highlighted[0]["column_letter"] == "DP"
    assert highlighted[0]["value"] == "new-budget"
    assert highlighted[0]["histories"][0]["before"] == "old-budget"


def test_word_diff_reports_paragraph_and_table_changes(tmp_path):
    left = tmp_path / "left.docx"
    right = tmp_path / "right.docx"
    _write_word(left, "본문 버전 A", "Alpha")
    _write_word(right, "본문 버전 B", "Beta")

    result = run_consistency_check(
        [
            {'id': 1, 'path': str(left), 'name': 'left.docx', 'file_type': 'Word'},
            {'id': 2, 'path': str(right), 'name': 'right.docx', 'file_type': 'Word'},
        ]
    )

    assert result["mode"] == "word"
    assert any(change["change_type"] == "replace" for change in result["word"]["changes"])
    before_types = {
        block["block_type"]
        for change in result["word"]["changes"]
        for block in change["before"]
    }
    assert {"paragraph", "table_row"} <= before_types


def test_word_diff_reports_best_effort_page_numbers(tmp_path):
    left = tmp_path / "left-page.docx"
    right = tmp_path / "right-page.docx"
    _write_word_with_second_page_change(left, "두 번째 페이지 원본")
    _write_word_with_second_page_change(right, "두 번째 페이지 수정")

    result = run_consistency_check(
        [
            {'id': 1, 'path': str(left), 'name': 'left-page.docx', 'file_type': 'Word'},
            {'id': 2, 'path': str(right), 'name': 'right-page.docx', 'file_type': 'Word'},
        ]
    )

    change = next(change for change in result["word"]["changes"] if change["change_type"] == "replace")
    assert change["before"][0]["page_number"] == 2
    assert change["after"][0]["page_number"] == 2


def test_ppt_diff_reports_slide_insert_and_update(tmp_path):
    left = tmp_path / "left.pptx"
    right = tmp_path / "right.pptx"
    _write_ppt(left, first_body="원본 본문", include_inserted_slide=False)
    _write_ppt(right, first_body="수정된 본문", include_inserted_slide=True)

    result = run_consistency_check(
        [
            {'id': 1, 'path': str(left), 'name': 'left.pptx', 'file_type': 'PowerPoint'},
            {'id': 2, 'path': str(right), 'name': 'right.pptx', 'file_type': 'PowerPoint'},
        ]
    )

    assert result["mode"] == "ppt"
    change_types = {change["change_type"] for change in result["ppt"]["changes"]}
    assert {"slide_insert", "slide_update"} <= change_types


def test_check_api_rejects_mixed_file_types(monkeypatch):
    rows = {
        1: {'id': 1, 'path': '/tmp/a.xlsx', 'name': 'a.xlsx', 'file_type': 'Excel'},
        2: {'id': 2, 'path': '/tmp/b.docx', 'name': 'b.docx', 'file_type': 'Word'},
    }
    monkeypatch.setattr("backend.api.check.get_file_by_id", lambda file_id: rows[file_id])
    monkeypatch.setattr("backend.api.check.os.path.exists", lambda _: True)

    with pytest.raises(HTTPException) as exc_info:
        consistency_check(CheckRequest(file_ids=[1, 2]))

    assert exc_info.value.status_code == 400


def test_check_api_reuses_comparison_cache_until_index_mtime_changes(tmp_path, monkeypatch):
    from backend.database import init_db, register_file, update_file_mtime

    monkeypatch.setattr("backend.database.DB_PATH", tmp_path / "test.db")
    monkeypatch.setattr("backend.database.DB_DIR", tmp_path)
    init_db()

    left = tmp_path / "left.docx"
    right = tmp_path / "right.docx"
    left.write_text("left", encoding="utf-8")
    right.write_text("right", encoding="utf-8")
    left_id = register_file(str(left), 'left.docx', 'Word', 0)
    right_id = register_file(str(right), 'right.docx', 'Word', 0)
    update_file_mtime(left_id, left.stat().st_mtime)
    update_file_mtime(right_id, right.stat().st_mtime)

    calls = 0

    def fake_run(file_infos):
        nonlocal calls
        calls += 1
        return {
            "mode": "word",
            "excel": None,
            "word": {
                "files": [
                    {"file_id": file_infos[0]["id"], "file_name": file_infos[0]["name"]},
                    {"file_id": file_infos[1]["id"], "file_name": file_infos[1]["name"]},
                ],
                "changes": [],
            },
            "ppt": None,
        }

    monkeypatch.setattr("backend.api.check.run_consistency_check", fake_run)

    first = consistency_check(CheckRequest(file_ids=[left_id, right_id]))
    second = consistency_check(CheckRequest(file_ids=[left_id, right_id]))

    assert first.mode == second.mode == "word"
    assert calls == 1

    right.write_text("right changed", encoding="utf-8")
    newer_mtime = right.stat().st_mtime + 5
    os.utime(right, (newer_mtime, newer_mtime))
    third = consistency_check(CheckRequest(file_ids=[left_id, right_id]))

    assert third.mode == "word"
    assert calls == 1
    assert any(warning.type == "source_may_be_newer" for warning in third.metadata.warnings)

    update_file_mtime(right_id, right.stat().st_mtime)
    fourth = consistency_check(CheckRequest(file_ids=[left_id, right_id]))

    assert fourth.mode == "word"
    assert calls == 2


def test_check_api_cache_recovers_from_corrupt_entry(tmp_path, monkeypatch):
    from backend.api.check import _comparison_cache_key
    from backend.database import init_db, register_file

    monkeypatch.setattr("backend.database.DB_PATH", tmp_path / "test.db")
    monkeypatch.setattr("backend.database.DB_DIR", tmp_path)
    init_db()

    left = tmp_path / "left.xlsx"
    right = tmp_path / "right.xlsx"
    left.write_text("left", encoding="utf-8")
    right.write_text("right", encoding="utf-8")
    left_id = register_file(str(left), 'left.xlsx', 'Excel', 1)
    right_id = register_file(str(right), 'right.xlsx', 'Excel', 1)

    calls = 0

    def fake_run(file_infos):
        nonlocal calls
        calls += 1
        return {
            "mode": "excel",
            "excel": {"total_keys": 0, "matched_keys": 0, "issues": []},
            "word": None,
            "ppt": None,
        }

    monkeypatch.setattr("backend.api.check.run_consistency_check", fake_run)

    consistency_check(CheckRequest(file_ids=[left_id, right_id]))
    consistency_check(CheckRequest(file_ids=[left_id, right_id]))

    assert calls == 1

    file_infos = [
        {
            "id": left_id,
            "path": str(left),
            "name": "left.xlsx",
            "file_type": "Excel",
        },
        {
            "id": right_id,
            "path": str(right),
            "name": "right.xlsx",
            "file_type": "Excel",
        },
    ]
    cache_key = _comparison_cache_key(file_infos, "version_history")
    import sqlite3

    conn = sqlite3.connect(tmp_path / "test.db")
    conn.execute("UPDATE comparison_cache SET result_json = ? WHERE cache_key = ?", ("{bad json", cache_key))
    conn.commit()
    conn.close()

    consistency_check(CheckRequest(file_ids=[left_id, right_id]))

    assert calls == 2
