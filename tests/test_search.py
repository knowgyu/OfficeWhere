import os
import sqlite3
import time
import tempfile
from datetime import datetime

import pytest
from openpyxl import Workbook

from backend.core.search_cache import reset_search_cache_for_tests
from backend.core.indexer import index_file, inspect_and_chunk, reindex_all, search, _sanitize_fts_query
from backend.database import (
    init_db,
    register_file,
    delete_file,
    get_all_files,
    search_chunks,
    save_file_chunks,
    update_file_mtime,
)


@pytest.fixture(autouse=True)
def setup_db(tmp_path, monkeypatch):
    db_path = tmp_path / "test.db"
    monkeypatch.setattr("backend.database.DB_PATH", db_path)
    monkeypatch.setattr("backend.database.DB_DIR", tmp_path)
    init_db()
    reset_search_cache_for_tests()


def _make_excel(path: str, data: dict):
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Sheet1"
    headers = list(data.keys())
    worksheet.append(headers)
    row_count = max((len(values) for values in data.values()), default=0)
    for row_index in range(row_count):
        worksheet.append([
            values[row_index] if row_index < len(values) else ""
            for values in data.values()
        ])
    workbook.save(path)


def _escape_pdf_text(value: str) -> str:
    return value.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def _make_pdf(path: str, pages: list[str]):
    """Write a minimal text PDF fixture without adding another test dependency."""

    objects: list[str] = [
        "<< /Type /Catalog /Pages 2 0 R >>",
        "",
        "<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    page_object_ids: list[int] = []
    for text in pages:
        page_object_id = len(objects) + 1
        content_object_id = page_object_id + 1
        page_object_ids.append(page_object_id)
        stream = f"BT /F1 14 Tf 72 720 Td ({_escape_pdf_text(text)}) Tj ET\n"
        objects.append(
            f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            f"/Resources << /Font << /F1 3 0 R >> >> /Contents {content_object_id} 0 R >>"
        )
        objects.append(f"<< /Length {len(stream.encode('latin-1'))} >>\nstream\n{stream}endstream")

    objects[1] = (
        f"<< /Type /Pages /Kids [{' '.join(f'{object_id} 0 R' for object_id in page_object_ids)}] "
        f"/Count {len(page_object_ids)} >>"
    )

    output = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for object_id, body in enumerate(objects, start=1):
        offsets.append(len(output))
        output.extend(f"{object_id} 0 obj\n{body}\nendobj\n".encode("latin-1"))
    xref_offset = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n0000000000 65535 f \n".encode("latin-1"))
    for offset in offsets[1:]:
        output.extend(f"{offset:010d} 00000 n \n".encode("latin-1"))
    output.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
        f"startxref\n{xref_offset}\n%%EOF\n".encode("latin-1")
    )
    with open(path, "wb") as handle:
        handle.write(output)


def test_index_and_search_excel(tmp_path):
    xlsx = str(tmp_path / "sample.xlsx")
    _make_excel(xlsx, {"과제명": ["DFBA 챗봇", "스마트팜"], "담당자": ["홍길동", "김철수"]})

    file_id = register_file(xlsx, 'sample.xlsx', 'xlsx', 2)
    chunk_count = index_file(file_id, xlsx)

    assert chunk_count > 0

    results = search("DFBA")
    assert len(results) > 0
    assert any("DFBA" in r["snippet"] or "DFBA" in r.get("content", "") for r in results)



def test_reindex_all_prunes_legacy_text_rows(tmp_path):
    txt = tmp_path / "legacy.txt"
    txt.write_text("old text", encoding="utf-8")
    register_file(str(txt), 'legacy.txt', 'Text', 0)

    stats = reindex_all()

    assert stats == {"success": 0, "failed": 0, "skipped": 1}
    assert get_all_files() == []
    assert txt.exists()


def test_reindex_all_invalidates_search_response_cache_epoch(tmp_path):
    from backend.core.search_cache import current_epoch

    before = current_epoch()

    reindex_all()

    assert current_epoch() > before


def test_search_returns_location(tmp_path):
    xlsx = str(tmp_path / "loc.xlsx")
    _make_excel(xlsx, {"항목": ["알파", "베타"], "값": ["100", "200"]})

    file_id = register_file(xlsx, 'loc.xlsx', 'xlsx', 2)
    index_file(file_id, xlsx)

    results = search("알파")
    assert len(results) > 0
    assert results[0]["location"] == "Sheet1 시트 | 2행 A열"


def test_index_and_search_excel_all_visible_sheets(tmp_path):
    xlsx = tmp_path / "multi-sheet.xlsx"
    workbook = Workbook()
    workbook.active.title = "요약"
    workbook.active["A1"] = "요약내용"
    detail = workbook.create_sheet("세부")
    detail["C3"] = "세컨드시트키워드"
    workbook.save(xlsx)

    file_id = register_file(str(xlsx), 'multi-sheet.xlsx', 'Excel', 0)
    chunk_count = index_file(file_id, str(xlsx))

    assert chunk_count == 2
    results = search("세컨드시트키워드")
    assert any(result["location"] == "세부 시트 | 3행 C열" for result in results)


def test_search_excel_header_uses_cell_location(tmp_path):
    xlsx = str(tmp_path / "header.xlsx")
    _make_excel(xlsx, {"항목": ["알파"], "담당자": ["홍길동"]})

    file_id = register_file(xlsx, 'header.xlsx', 'xlsx', 2)
    index_file(file_id, xlsx)

    results = search("담당자")

    assert len(results) > 0
    assert results[0]["location"] == "Sheet1 시트 | 1행 B열"


def test_index_excel_uses_used_range(tmp_path):
    xlsx = str(tmp_path / "stale.xlsx")
    _make_excel(xlsx, {"항목": ["알파"], "새열": ["범위밖키워드"]})
    file_id = register_file(xlsx, 'stale.xlsx', 'Excel', 1)
    chunk_count = index_file(file_id, xlsx)

    assert chunk_count > 0
    results = search("범위밖키워드")
    assert len(results) == 1
    assert results[0]["location"] == "Sheet1 시트 | 2행 B열"


def test_inspect_and_chunk_indexes_used_range(tmp_path):
    xlsx = str(tmp_path / "recover.xlsx")
    _make_excel(xlsx, {"항목": ["알파"], "새열": ["복구키워드"]})
    info, chunks = inspect_and_chunk(xlsx)

    assert info["columns"] == ["항목", "새열"]
    assert any(chunk["content"] == "복구키워드" for chunk in chunks)


def test_search_no_results_for_missing_term(tmp_path):
    xlsx = str(tmp_path / "empty.xlsx")
    _make_excel(xlsx, {"항목": ["ABC"], "값": ["123"]})

    file_id = register_file(xlsx, 'empty.xlsx', 'xlsx', 2)
    index_file(file_id, xlsx)

    results = search("존재하지않는단어XYZ")
    assert results == []


def test_search_matches_korean_substrings_inside_words(tmp_path):
    doc_path = tmp_path / "meeting.docx"

    file_id = register_file(str(doc_path), 'meeting.docx', 'Word', 0)
    save_file_chunks(file_id, [{"location": "문단", "content": "주간 회의록 작성 후 공유"}])

    results = search("회의")

    assert len(results) == 1
    assert results[0]["file_id"] == file_id
    assert "**회의**" in results[0]["snippet"]


def test_search_no_longer_guarantees_hangul_choseong(tmp_path):
    doc_path = tmp_path / "meeting.docx"

    file_id = register_file(str(doc_path), 'meeting.docx', 'Word', 0)
    save_file_chunks(file_id, [{"location": "문단", "content": "주간 회의록 작성 후 공유"}])

    results = search("ㅎㅇㄹ")

    assert results == []


def test_search_matches_long_korean_substring_with_fast_path(tmp_path):
    doc_path = tmp_path / "project.docx"

    file_id = register_file(str(doc_path), 'project.docx', 'Word', 0)
    save_file_chunks(file_id, [{"location": "문단", "content": "프로젝트 상태 보고서"}])

    results = search("프로젝")

    assert len(results) == 1
    assert results[0]["file_id"] == file_id
    assert "**프로젝**" in results[0]["snippet"]


def test_search_documents_degrades_content_search_while_search_index_repairs(tmp_path):
    from backend.application.search_service import search_documents
    from backend.database import set_setting
    from backend.models.schemas import SearchRequest

    file_id = register_file(str(tmp_path / "meeting.docx"), "meeting.docx", "Word", 0)
    save_file_chunks(file_id, [{"location": "문단", "content": "주간 회의록"}])
    set_setting("search_index_version", "4")
    init_db()

    content_response = search_documents(SearchRequest(query="회의", search_scope="content"))

    assert content_response.results == []
    assert content_response.search_index_stale is True
    assert content_response.search_index_state == "repair_needed"

    filename_response = search_documents(SearchRequest(query="meeting", search_scope="filename_content"))

    assert filename_response.search_index_stale is True
    assert [item.location for item in filename_response.results] == ["파일명"]
    assert filename_response.results[0].file_id == file_id


def test_init_db_removes_unused_base_file_search():
    from backend.database import DB_PATH

    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute("SELECT name FROM sqlite_master WHERE name = 'file_search'")
    assert cursor.fetchone() is None
    cursor.execute("SELECT name FROM sqlite_master WHERE type = 'trigger' AND name IN ('chunks_ai', 'chunks_ad')")
    assert cursor.fetchall() == []
    cursor.execute("SELECT name FROM sqlite_master WHERE name = 'file_search_ko'")
    assert cursor.fetchone() == ("file_search_ko",)
    cursor.execute(
        """
        SELECT name FROM sqlite_master
        WHERE name IN ('file_search_ko_docsize', 'file_search_trigram_docsize')
        """
    )
    assert cursor.fetchall() == []
    conn.close()


def test_init_db_migrates_legacy_base_file_search(tmp_path, monkeypatch):
    from backend.database import DB_PATH, init_db, register_file, save_file_chunks

    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute(
        """
        CREATE VIRTUAL TABLE file_search USING fts5(
            content,
            content='file_chunks',
            content_rowid='id',
            tokenize='unicode61'
        )
        """
    )
    cursor.execute(
        """
        CREATE TRIGGER chunks_ai AFTER INSERT ON file_chunks BEGIN
            INSERT INTO file_search(rowid, content) VALUES (new.id, new.content);
        END
        """
    )
    cursor.execute(
        """
        CREATE TRIGGER chunks_ad AFTER DELETE ON file_chunks BEGIN
            INSERT INTO file_search(file_search, rowid, content)
            VALUES ('delete', old.id, old.content);
        END
        """
    )
    conn.commit()
    conn.close()

    init_db()

    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute("SELECT name FROM sqlite_master WHERE name = 'file_search'")
    assert cursor.fetchone() is None
    cursor.execute("SELECT name FROM sqlite_master WHERE type = 'trigger' AND name IN ('chunks_ai', 'chunks_ad')")
    assert cursor.fetchall() == []
    conn.close()

    file_id = register_file(str(tmp_path / 'meeting.docx'), 'meeting.docx', 'Word', 0)
    save_file_chunks(file_id, [{"location": "문단", "content": "주간 회의록"}])

    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute("SELECT rowid FROM file_search_ko WHERE file_search_ko MATCH ?", ('"회의"',))
    assert cursor.fetchone() is not None
    conn.close()


def test_init_db_defers_search_index_repair_until_background_repair(tmp_path, monkeypatch):
    from backend.database import DB_PATH, get_search_index_status, repair_search_indexes, set_setting

    file_id = register_file(str(tmp_path / 'meeting.docx'), 'meeting.docx', 'Word', 0)
    save_file_chunks(file_id, [{"location": "문단", "content": "주간 회의록"}])

    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute("DROP TRIGGER IF EXISTS chunks_ai_ko")
    cursor.execute("DROP TRIGGER IF EXISTS chunks_ad_ko")
    cursor.execute("DROP TRIGGER IF EXISTS chunks_ai_trigram")
    cursor.execute("DROP TRIGGER IF EXISTS chunks_ad_trigram")
    cursor.execute("DROP TABLE IF EXISTS file_search_ko")
    cursor.execute("DROP TABLE IF EXISTS file_search_trigram")
    cursor.execute(
        """
        CREATE VIRTUAL TABLE file_search_ko USING fts5(
            search_text,
            content='file_chunks',
            content_rowid='id',
            tokenize='unicode61'
        )
        """
    )
    cursor.execute("INSERT INTO file_search_ko(file_search_ko) VALUES ('rebuild')")
    conn.commit()
    conn.close()
    set_setting("search_index_version", "4")

    with monkeypatch.context() as patch:
        patch.setattr(
            "backend.database._rebuild_search_indexes",
            lambda *_args, **_kwargs: (_ for _ in ()).throw(AssertionError("init_db must not rebuild FTS")),
        )
        patch.setattr(
            "backend.database._refresh_search_text",
            lambda *_args, **_kwargs: (_ for _ in ()).throw(AssertionError("init_db must not refresh all chunks")),
        )
        init_db()

    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute(
        """
        SELECT name FROM sqlite_master
        WHERE name IN ('file_search_ko_docsize', 'file_search_trigram_docsize')
        """
    )
    assert cursor.fetchall() == [("file_search_ko_docsize",)]
    conn.close()
    status = get_search_index_status()
    assert status["state"] == "repair_needed"
    assert status["stale"] is True

    repair_search_indexes(reason="test")

    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute(
        """
        SELECT name FROM sqlite_master
        WHERE name IN ('file_search_ko_docsize', 'file_search_trigram_docsize')
        """
    )
    assert cursor.fetchall() == []
    cursor.execute("SELECT rowid FROM file_search_ko WHERE file_search_ko MATCH ?", ('"회의"',))
    assert cursor.fetchone() is not None
    conn.close()
    status = get_search_index_status()
    assert status["state"] == "ready"
    assert status["stale"] is False


def test_reindex_on_file_change(tmp_path):
    xlsx = str(tmp_path / "change.xlsx")
    _make_excel(xlsx, {"항목": ["원래값"], "값": ["1"]})

    file_id = register_file(xlsx, 'change.xlsx', 'xlsx', 1)
    index_file(file_id, xlsx)

    assert search("원래값") != []
    assert search("수정값") == []

    # 파일 수정 후 재인덱싱
    _make_excel(xlsx, {"항목": ["수정값"], "값": ["2"]})
    index_file(file_id, xlsx)

    assert search("수정값") != []


def test_sanitize_fts_query():
    assert _sanitize_fts_query("hello world") == '"hello world"'
    assert _sanitize_fts_query('foo "bar"') == '"foo bar"'
    assert _sanitize_fts_query("  ") == '""'


def test_spaced_query_prefers_exact_phrase_matches(tmp_path):
    from backend.api.search import search_files
    from backend.models.schemas import SearchRequest

    phrase_id = register_file(str(tmp_path / 'phrase.docx'), 'phrase.docx', 'Word', 0)
    loose_id = register_file(str(tmp_path / 'loose.docx'), 'loose.docx', 'Word', 0)
    save_file_chunks(phrase_id, [{"location": "문단", "content": "데이터 연계 전략을 정리했습니다"}])
    save_file_chunks(loose_id, [{"location": "문단", "content": "데이터 기반 보고서와 별도 연계표"}])

    response = search_files(SearchRequest(query="데이터 연계", search_scope="content"))

    assert {item.file_id for item in response.results} == {phrase_id}
    assert "**데이터 연계**" in response.results[0].snippet


def test_index_performance_excel(tmp_path):
    """500행 Excel 인덱싱이 5초 내 완료."""
    xlsx = str(tmp_path / "perf.xlsx")
    _make_excel(xlsx, {
        "과제명": [f"과제_{i}" for i in range(500)],
        "담당자": [f"담당자_{i}" for i in range(500)],
        "예산": [str(i * 1000) for i in range(500)],
    })

    file_id = register_file(xlsx, 'perf.xlsx', 'xlsx', 3)
    start = time.perf_counter()
    chunk_count = index_file(file_id, xlsx)
    elapsed = time.perf_counter() - start

    assert chunk_count > 0
    assert elapsed < 5.0, f"인덱싱 {elapsed:.2f}초 — 5초 초과"


def test_search_performance(tmp_path):
    """1000청크 인덱싱 후 검색이 0.5초 내 완료."""
    xlsx = str(tmp_path / "search_perf.xlsx")
    _make_excel(xlsx, {
        "항목": [f"항목_{i}" for i in range(1000)],
        "내용": [f"내용 데이터 {i} 테스트" for i in range(1000)],
    })

    file_id = register_file(xlsx, 'search_perf.xlsx', 'xlsx', 2)
    index_file(file_id, xlsx)

    start = time.perf_counter()
    results = search("데이터")
    elapsed = time.perf_counter() - start

    assert elapsed < 0.5, f"검색 {elapsed:.3f}초 — 0.5초 초과"


def test_search_chunks_filters_file_type(tmp_path):
    from backend.database import save_file_chunks

    doc_id = register_file(str(tmp_path / 'note.docx'), 'note.docx', 'Word', 0)
    ppt_id = register_file(str(tmp_path / 'note.pptx'), 'note.pptx', 'PowerPoint', 0)
    save_file_chunks(doc_id, [{"location": "문단", "content": "공통 키워드"}])
    save_file_chunks(ppt_id, [{"location": "슬라이드 1", "content": "공통 키워드"}])

    word_results = search_chunks('"공통"', file_types=["Word"])

    assert {item["file_type"] for item in word_results} == {"Word"}


def test_search_api_filters_filename_and_content(tmp_path):
    from backend.api.search import search_files
    from backend.database import save_file_chunks
    from backend.models.schemas import SearchRequest

    doc_id = register_file(str(tmp_path / 'alpha.docx'), 'alpha.docx', 'Word', 0)
    ppt_id = register_file(str(tmp_path / 'alpha.pptx'), 'alpha.pptx', 'PowerPoint', 0)
    save_file_chunks(doc_id, [{"location": "문단", "content": "검색 대상"}])
    save_file_chunks(ppt_id, [{"location": "슬라이드 1", "content": "검색 대상"}])

    response = search_files(SearchRequest(query="alpha", file_types=["pptx"]))

    assert response.total == 1
    assert response.results[0].file_type == "PowerPoint"


def test_search_api_default_scope_includes_filename_and_content(tmp_path):
    from backend.api.search import search_files
    from backend.database import save_file_chunks
    from backend.models.schemas import SearchRequest

    filename_id = register_file(str(tmp_path / 'scope-key.docx'), 'scope-key.docx', 'Word', 0)
    content_id = register_file(str(tmp_path / 'content.docx'), 'content.docx', 'Word', 0)
    save_file_chunks(filename_id, [{"location": "문단", "content": "다른 내용"}])
    save_file_chunks(content_id, [{"location": "문단", "content": "scope-key 본문"}])

    response = search_files(SearchRequest(query="scope-key"))

    assert response.total == 2
    assert {item.file_id for item in response.results} == {filename_id, content_id}


def test_search_api_filename_scope_excludes_content_only_matches(tmp_path):
    from backend.api.search import search_files
    from backend.database import save_file_chunks
    from backend.models.schemas import SearchRequest

    filename_id = register_file(str(tmp_path / 'onlyname.docx'), 'onlyname.docx', 'Word', 0)
    content_id = register_file(str(tmp_path / 'body.docx'), 'body.docx', 'Word', 0)
    save_file_chunks(filename_id, [{"location": "문단", "content": "다른 내용"}])
    save_file_chunks(content_id, [{"location": "문단", "content": "onlyname 본문"}])

    response = search_files(SearchRequest(query="onlyname", search_scope="filename"))

    assert response.total == 1
    assert response.results[0].file_id == filename_id
    assert response.results[0].location == "파일명"
    assert "**onlyname**" in response.results[0].snippet


def test_search_api_content_scope_excludes_filename_only_matches(tmp_path):
    from backend.api.search import search_files
    from backend.database import save_file_chunks
    from backend.models.schemas import SearchRequest

    filename_id = register_file(str(tmp_path / 'bodyonly.docx'), 'bodyonly.docx', 'Word', 0)
    content_id = register_file(str(tmp_path / 'note.docx'), 'note.docx', 'Word', 0)
    save_file_chunks(filename_id, [{"location": "문단", "content": "다른 내용"}])
    save_file_chunks(content_id, [{"location": "문단", "content": "bodyonly 본문"}])

    response = search_files(SearchRequest(query="bodyonly", search_scope="content"))

    assert response.total == 1
    assert response.results[0].file_id == content_id
    assert response.results[0].location == "문단"


def test_search_api_filename_scope_respects_file_type_filter(tmp_path):
    from backend.api.search import search_files
    from backend.models.schemas import SearchRequest

    register_file(str(tmp_path / 'scope.docx'), 'scope.docx', 'Word', 0)
    ppt_id = register_file(str(tmp_path / 'scope.pptx'), 'scope.pptx', 'PowerPoint', 0)

    response = search_files(SearchRequest(query="scope", file_types=["pptx"], search_scope="filename"))

    assert response.total == 1
    assert response.results[0].file_id == ppt_id
    assert response.results[0].file_type == "PowerPoint"


def test_search_request_logs_request_level_timings_and_fallback_reason(tmp_path, monkeypatch):
    from backend.api.search import search_files
    from backend.models.schemas import SearchRequest

    events = []
    monkeypatch.setattr(
        "backend.application.search_service.log_index_perf",
        lambda event, **fields: events.append((event, fields)),
    )
    file_id = register_file(str(tmp_path / "log-target.docx"), "log-target.docx", "Word", 0)
    update_file_mtime(file_id, 1_700_000_000.0)

    response = search_files(SearchRequest(query="log-target", search_scope="filename"))

    assert response.results[0].file_id == file_id
    event, fields = events[-1]
    assert event == "search_request_done"
    assert fields["success"] is True
    assert fields["request_id"]
    assert fields["search_scope"] == "filename"
    assert fields["cache_status"] == "miss"
    assert fields["filename_source"] == "db_like"
    assert fields["filename_fallback_reason"] in {"non_windows", "disabled", "sdk_dll_missing"}
    assert fields["filename_ms"] >= 0
    assert fields["content_ms"] == 0
    assert fields["merge_ms"] >= 0
    assert fields["row_count"] == 1


def test_search_response_cache_hits_then_invalidates_on_mtime_update(tmp_path, monkeypatch):
    from backend.api.search import search_files
    from backend.models.schemas import SearchRequest

    events = []
    monkeypatch.setattr(
        "backend.application.search_service.log_index_perf",
        lambda event, **fields: events.append((event, fields)),
    )
    file_id = register_file(str(tmp_path / "cache-target.docx"), "cache-target.docx", "Word", 0)
    update_file_mtime(file_id, 1_700_000_000.0)
    req = SearchRequest(query="cache-target", search_scope="filename")

    first = search_files(req)
    second = search_files(req)
    update_file_mtime(file_id, 1_700_000_111.0)
    third = search_files(req)

    assert first.results[0].file_mtime == 1_700_000_000.0
    assert second.results[0].file_mtime == 1_700_000_000.0
    assert third.results[0].file_mtime == 1_700_000_111.0
    statuses = [fields["cache_status"] for event, fields in events if event == "search_request_done"]
    assert statuses == ["miss", "hit", "miss"]


def test_search_storage_prewarm_touches_search_tables_without_writes(monkeypatch):
    from backend.database import prewarm_search_storage

    events = []
    monkeypatch.setattr(
        "backend.database.log_index_perf",
        lambda event, **fields: events.append((event, fields)),
    )

    result = prewarm_search_storage()

    assert result["success"] is True
    assert result["probe_count"] == 6
    assert "file_search_ko" in result["probes"]
    assert "file_search_trigram" in result["probes"]
    assert events[-1][0] == "search_storage_prewarm_done"
    assert events[-1][1]["success"] is True


def test_filename_content_search_cache_miss_reuses_one_read_connection(tmp_path, monkeypatch):
    from backend import database
    from backend.api.search import search_files
    from backend.core.everything_scanner import EverythingDiscovery
    from backend.models.schemas import SearchRequest

    events = []
    monkeypatch.setattr(
        "backend.application.search_service.log_index_perf",
        lambda event, **fields: events.append((event, fields)),
    )
    monkeypatch.setattr(
        "backend.application.search_service.discover_filename_candidates",
        lambda _query: EverythingDiscovery(unavailable_reason="disabled"),
    )

    file_id = register_file(
        str(tmp_path / "single-open-target.docx"),
        "single-open-target.docx",
        "Word",
        0,
    )
    save_file_chunks(file_id, [{"location": "문단", "content": "single-open-target 본문"}])
    reset_search_cache_for_tests()

    original_connect = database._connect
    connect_count = 0

    def counted_connect():
        nonlocal connect_count
        connect_count += 1
        return original_connect()

    monkeypatch.setattr(database, "_connect", counted_connect)

    response = search_files(SearchRequest(query="single-open-target", search_scope="filename_content"))

    request_events = [fields for event, fields in events if event == "search_request_done"]
    assert request_events[-1]["cache_status"] == "miss"
    assert connect_count == 1
    assert response.file_count == 1
    assert {item.location for item in response.results} >= {"파일명", "문단"}


def test_everything_filename_content_path_reuses_one_read_connection(tmp_path, monkeypatch):
    from backend import database
    from backend.api.search import search_files
    from backend.core.everything_scanner import EverythingDiscovery
    from backend.models.schemas import SearchRequest

    events = []
    monkeypatch.setattr(
        "backend.application.search_service.log_index_perf",
        lambda event, **fields: events.append((event, fields)),
    )

    older_id = register_file(str(tmp_path / "everything-open-a.docx"), "everything-open-a.docx", "Word", 0)
    newer_id = register_file(str(tmp_path / "everything-open-b.docx"), "everything-open-b.docx", "Word", 0)
    save_file_chunks(older_id, [{"location": "문단", "content": "everything-open 본문 A"}])
    save_file_chunks(newer_id, [{"location": "문단", "content": "everything-open 본문 B"}])
    unregistered = tmp_path / "everything-open-unregistered.docx"
    monkeypatch.setattr(
        "backend.application.search_service.discover_filename_candidates",
        lambda _query: EverythingDiscovery(
            paths=[
                str(tmp_path / "everything-open-a.docx"),
                str(tmp_path / "everything-open-b.docx"),
                str(unregistered),
            ],
            queried_count=3,
        ),
    )
    reset_search_cache_for_tests()

    original_connect = database._connect
    connect_count = 0

    def counted_connect():
        nonlocal connect_count
        connect_count += 1
        return original_connect()

    monkeypatch.setattr(database, "_connect", counted_connect)

    response = search_files(SearchRequest(query="everything-open", search_scope="filename_content", file_limit=2))

    request_events = [fields for event, fields in events if event == "search_request_done"]
    assert request_events[-1]["cache_status"] == "miss"
    assert request_events[-1]["filename_source"] == "everything_sdk"
    assert connect_count == 1
    assert response.file_count == 2
    assert response.results[0].file_id == newer_id
    assert all("unregistered" not in item.path for item in response.results)


def test_borrowed_search_connection_remains_open_and_row_compatible(tmp_path):
    from backend import database

    file_id = register_file(str(tmp_path / "borrowed-open.docx"), "borrowed-open.docx", "Word", 0)
    save_file_chunks(file_id, [{"location": "문단", "content": "borrowed-open 본문"}])

    conn = database._connect()
    original_row_factory = conn.row_factory
    try:
        status = database.get_search_index_status(conn=conn)
        assert conn.row_factory is original_row_factory

        filename_rows = database.search_file_names("borrowed-open", conn=conn)
        assert conn.row_factory is original_row_factory

        path_rows = database.search_file_names_by_paths(
            "borrowed-open",
            [str(tmp_path / "borrowed-open.docx")],
            conn=conn,
        )
        assert conn.row_factory is original_row_factory

        content_rows = database.search_chunks('"borrowed-open"', raw_query="borrowed-open", conn=conn)
        assert conn.row_factory is original_row_factory

        assert status["state"] == "ready"
        assert filename_rows[0]["id"] == file_id
        assert path_rows[0]["id"] == file_id
        assert content_rows[0]["file_id"] == file_id
        assert conn.execute("SELECT 1").fetchone()[0] == 1
    finally:
        conn.close()


def test_everything_filename_candidates_preserve_db_authority_order_and_pagination(tmp_path, monkeypatch):
    from backend.api.search import search_files
    from backend.core.everything_scanner import EverythingDiscovery
    from backend.models.schemas import SearchRequest

    older_id = register_file(str(tmp_path / "report-a.docx"), "report-a.docx", "Word", 0)
    newer_id = register_file(str(tmp_path / "report-b.pptx"), "report-b.pptx", "PowerPoint", 0)
    register_file(str(tmp_path / "report-c.xlsx"), "report-c.xlsx", "Excel", 0)
    unregistered = tmp_path / "report-unregistered.pdf"
    monkeypatch.setattr(
        "backend.application.search_service.discover_filename_candidates",
        lambda _query: EverythingDiscovery(
            paths=[
                str(tmp_path / "report-a.docx"),
                str(tmp_path / "report-b.pptx"),
                str(unregistered),
            ],
            queried_count=3,
        ),
    )

    response = search_files(SearchRequest(query="report", search_scope="filename", file_limit=1))

    assert [item.file_id for item in response.results] == [newer_id]
    assert response.has_more is True
    assert older_id != newer_id
    assert all("unregistered" not in item.path for item in response.results)


def test_everything_filename_candidate_fallback_keeps_db_results(tmp_path, monkeypatch):
    from backend.api.search import search_files
    from backend.core.everything_scanner import EverythingDiscovery
    from backend.models.schemas import SearchRequest

    file_id = register_file(str(tmp_path / "fallback-report.docx"), "fallback-report.docx", "Word", 0)
    monkeypatch.setattr(
        "backend.application.search_service.discover_filename_candidates",
        lambda _query: EverythingDiscovery(unavailable_reason="candidate_limit_exceeded", queried_count=999),
    )

    response = search_files(SearchRequest(query="fallback-report", search_scope="filename"))

    assert [item.file_id for item in response.results] == [file_id]


def test_search_api_content_scope_supports_excel_filter(tmp_path):
    from backend.api.search import search_files
    from backend.database import save_file_chunks
    from backend.models.schemas import SearchRequest

    excel_id = register_file(str(tmp_path / 'sheet.xlsx'), 'sheet.xlsx', 'Excel', 0)
    word_id = register_file(str(tmp_path / 'sheet.docx'), 'sheet.docx', 'Word', 0)
    save_file_chunks(excel_id, [{"location": "Sheet1 행 1", "content": "엑셀검색 키워드"}])
    save_file_chunks(word_id, [{"location": "문단", "content": "엑셀검색 키워드"}])

    response = search_files(SearchRequest(query="엑셀검색", file_types=["xlsx"], search_scope="content"))

    assert response.total == 1
    assert response.results[0].file_id == excel_id
    assert response.results[0].file_type == "Excel"


def test_word_search_locations_use_page_labels(tmp_path):
    from docx import Document

    doc_path = tmp_path / "paged.docx"
    document = Document()
    document.add_paragraph("첫 페이지")
    document.add_page_break()
    document.add_paragraph("두 번째 페이지 검색키워드")
    document.save(doc_path)

    file_id = register_file(str(doc_path), 'paged.docx', 'Word', 0)
    index_file(file_id, str(doc_path))

    results = search("검색키워드")

    assert len(results) == 1
    assert results[0]["location"] == "쪽 2"


def test_ppt_search_locations_hide_shape_details(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    ppt_path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text = "프로젝트 검색키워드"
    presentation.save(ppt_path)

    file_id = register_file(str(ppt_path), 'slides.pptx', 'PowerPoint', 0)
    index_file(file_id, str(ppt_path))

    results = search("검색키워드")

    assert len(results) >= 1
    assert {result["location"] for result in results} == {"슬라이드 1"}
    assert all("shape" not in result["location"] for result in results)


def test_pdf_inspect_index_and_search(tmp_path):
    pdf_path = tmp_path / "report.pdf"
    _make_pdf(str(pdf_path), ["first page data linkage", "second page budget"])

    info, chunks = inspect_and_chunk(str(pdf_path))

    assert info["file_type"] == "PDF"
    assert info["columns"] == ["페이지", "내용"]
    assert info["sample"][0][0] == "쪽 1"
    assert chunks[0]["location"] == "쪽 1"

    file_id = register_file(str(pdf_path), "report.pdf", "PDF", len(info["columns"]))
    chunk_count = index_file(file_id, str(pdf_path))
    results = search("budget", file_types=["PDF"])

    assert chunk_count == 2
    assert len(results) == 1
    assert results[0]["file_id"] == file_id
    assert results[0]["file_type"] == "PDF"
    assert results[0]["location"] == "쪽 2"


def test_pdf_page_extraction_failure_fails_whole_index(monkeypatch, tmp_path):
    from backend.core import pdf_analysis

    class GoodPage:
        def get_textpage(self):
            class TextPage:
                def get_text_bounded(self, errors="replace"):
                    return "first page"

                def close(self):
                    pass

            return TextPage()

        def close(self):
            pass

    class BrokenPage:
        def get_textpage(self):
            raise RuntimeError("damaged content stream")

        def close(self):
            pass

    class FakeDocument:
        def __enter__(self):
            return self

        def __exit__(self, _exc_type, _exc, _tb):
            self.close()

        def __len__(self):
            return 2

        def __getitem__(self, index: int):
            return [GoodPage(), BrokenPage()][index]

        def close(self):
            pass

    class FakePdfium:
        class PdfiumError(RuntimeError):
            err_code = None

        @staticmethod
        def PdfDocument(_path: str, password=""):
            return FakeDocument()

    monkeypatch.setattr(pdf_analysis, "_load_pdfium", lambda: FakePdfium)

    with pytest.raises(ValueError, match="PDF 2쪽 텍스트를 읽을 수 없습니다"):
        inspect_and_chunk(str(tmp_path / "broken.pdf"))


def test_search_api_filters_pdf_type(tmp_path):
    from backend.api.search import search_files
    from backend.models.schemas import SearchRequest

    pdf_id = register_file(str(tmp_path / "note.pdf"), "note.pdf", "PDF", 2)
    word_id = register_file(str(tmp_path / "note.docx"), "note.docx", "Word", 0)
    save_file_chunks(pdf_id, [{"location": "쪽 1", "content": "공통 PDF 키워드"}])
    save_file_chunks(word_id, [{"location": "문단", "content": "공통 PDF 키워드"}])

    response = search_files(SearchRequest(query="공통", file_types=["pdf"], search_scope="content"))

    assert response.total == 1
    assert response.results[0].file_id == pdf_id
    assert response.results[0].file_type == "PDF"


def test_search_api_excludes_temporary_folder_paths_without_prefix_bleed(tmp_path):
    from backend.api.search import search_files
    from backend.models.schemas import SearchRequest

    hidden_dir = tmp_path / "docs" / "a"
    visible_prefix_dir = tmp_path / "docs" / "abc"
    visible_dir = tmp_path / "docs" / "other"
    hidden_id = register_file(str(hidden_dir / "hidden.docx"), "hidden.docx", "Word", 0)
    prefix_id = register_file(str(visible_prefix_dir / "prefix.docx"), "prefix.docx", "Word", 0)
    visible_id = register_file(str(visible_dir / "visible.docx"), "visible.docx", "Word", 0)
    for file_id in [hidden_id, prefix_id, visible_id]:
        save_file_chunks(file_id, [{"location": "문단", "content": "폴더제외 키워드"}])

    response = search_files(
        SearchRequest(
            query="폴더제외",
            search_scope="content",
            excluded_folder_paths=[str(hidden_dir)],
        )
    )

    assert {item.file_id for item in response.results} == {prefix_id, visible_id}


def test_search_api_excludes_windows_style_folder_paths_without_prefix_bleed(tmp_path):
    from backend.api.search import search_files
    from backend.models.schemas import SearchRequest

    hidden_id = register_file(r"C:\Docs\A\hidden.docx", "hidden.docx", "Word", 0)
    visible_prefix_id = register_file(r"C:\Docs\ABC\prefix.docx", "prefix.docx", "Word", 0)
    visible_id = register_file(r"C:\Docs\Other\visible.docx", "visible.docx", "Word", 0)
    for file_id in [hidden_id, visible_prefix_id, visible_id]:
        save_file_chunks(file_id, [{"location": "문단", "content": "윈도우 제외 키워드"}])

    response = search_files(
        SearchRequest(
            query="윈도우 제외",
            search_scope="content",
            excluded_folder_paths=[r"c:/docs/a"],
        )
    )

    assert {item.file_id for item in response.results} == {visible_prefix_id, visible_id}


def test_search_api_filters_by_file_modified_date(tmp_path):
    from backend.api.search import search_files
    from backend.database import save_file_chunks, update_file_mtime
    from backend.models.schemas import SearchRequest

    old_id = register_file(str(tmp_path / 'old.docx'), '회의_old.docx', 'Word', 0)
    new_id = register_file(str(tmp_path / 'new.docx'), '회의_new.docx', 'Word', 0)
    save_file_chunks(old_id, [{"location": "문단", "content": "회의 자료"}])
    save_file_chunks(new_id, [{"location": "문단", "content": "회의 자료"}])
    update_file_mtime(old_id, datetime(2026, 3, 15).timestamp())
    update_file_mtime(new_id, datetime(2026, 4, 15).timestamp())

    response = search_files(
        SearchRequest(
            query="회의",
            modified_from="2026-04-01",
            modified_to="2026-04-30",
        )
    )

    matched_ids = {item.file_id for item in response.results}
    assert matched_ids == {new_id}


def test_search_api_caps_results_by_file_first(tmp_path):
    from backend.api.search import search_files
    from backend.database import save_file_chunks
    from backend.models.schemas import SearchRequest

    for index in range(25):
        file_id = register_file(str(tmp_path / f'note-{index}.docx'), f'note-{index}.docx', 'Word', 0)
        save_file_chunks(file_id, [{"location": "문단", "content": f"프로젝트 공통키워드 {index}"}])

    response = search_files(SearchRequest(query="공통키워드", search_scope="content", file_limit=20))

    assert response.file_limit == 20
    assert response.file_count == 20
    assert len({item.file_id for item in response.results}) == 20
    assert response.has_more is True


def test_search_api_honors_file_offset_for_content_results(tmp_path):
    from backend.api.search import search_files
    from backend.database import save_file_chunks
    from backend.models.schemas import SearchRequest

    file_ids = []
    for index in range(3):
        file_id = register_file(str(tmp_path / f'offset-{index}.docx'), f'offset-{index}.docx', 'Word', 0)
        save_file_chunks(file_id, [{"location": "문단", "content": f"프로젝트 오프셋검색 {index}"}])
        update_file_mtime(file_id, float(index))
        file_ids.append(file_id)

    response = search_files(
        SearchRequest(query="오프셋검색", search_scope="content", file_limit=1, file_offset=1)
    )

    assert response.file_count == 1
    assert {item.file_id for item in response.results} == {file_ids[1]}
    assert response.has_more is True


def test_search_api_honors_file_offset_for_default_content_fallback(tmp_path):
    from backend.api.search import search_files
    from backend.database import save_file_chunks
    from backend.models.schemas import SearchRequest

    file_ids = []
    for index in range(3):
        file_id = register_file(str(tmp_path / f'fallback-{index}.docx'), f'fallback-{index}.docx', 'Word', 0)
        save_file_chunks(file_id, [{"location": "문단", "content": f"프로젝트 기본오프셋검색 {index}"}])
        update_file_mtime(file_id, float(index))
        file_ids.append(file_id)

    response = search_files(
        SearchRequest(query="기본오프셋검색", search_scope="filename_content", file_limit=1, file_offset=1)
    )

    assert response.file_count == 1
    assert {item.file_id for item in response.results} == {file_ids[1]}
    assert response.has_more is True


def test_search_api_stops_has_more_at_max_file_window(tmp_path):
    from backend.api.search import search_files
    from backend.database import save_file_chunks
    from backend.models.schemas import SearchRequest

    for index in range(105):
        file_id = register_file(str(tmp_path / f'cap-{index}.docx'), f'cap-{index}.docx', 'Word', 0)
        save_file_chunks(file_id, [{"location": "문단", "content": f"프로젝트 최대창검색 {index}"}])
        update_file_mtime(file_id, float(index))

    response = search_files(
        SearchRequest(query="최대창검색", search_scope="content", file_limit=20, file_offset=80)
    )

    assert response.file_count == 20
    assert response.has_more is False


def test_search_api_honors_file_offset_for_filename_results(tmp_path):
    from backend.api.search import search_files
    from backend.models.schemas import SearchRequest

    file_ids = []
    for index in range(3):
        file_id = register_file(
            str(tmp_path / f'오프셋파일-{index}.docx'),
            f'오프셋파일-{index}.docx',
            'Word',
            0,
        )
        update_file_mtime(file_id, float(index))
        file_ids.append(file_id)

    response = search_files(
        SearchRequest(query="오프셋파일", search_scope="filename", file_limit=1, file_offset=1)
    )

    assert response.file_count == 1
    assert {item.file_id for item in response.results} == {file_ids[1]}
    assert response.has_more is True


def test_search_api_honors_per_file_limit_for_common_terms(tmp_path):
    from backend.api.search import search_files
    from backend.database import save_file_chunks
    from backend.models.schemas import SearchRequest

    for file_index in range(4):
        file_id = register_file(
            str(tmp_path / f'common-{file_index}.docx'),
            f'common-{file_index}.docx',
            'Word',
            0,
        )
        save_file_chunks(
            file_id,
            [
                {
                    "location": f"문단 {chunk_index}",
                    "content": f"프로젝트 공통키워드 반복 {file_index}-{chunk_index}",
                }
                for chunk_index in range(5)
            ],
        )
        update_file_mtime(file_id, float(file_index))

    response = search_files(
        SearchRequest(query="공통키워드", search_scope="content", file_limit=3, per_file_limit=2)
    )

    counts_by_file: dict[int, int] = {}
    for item in response.results:
        counts_by_file[item.file_id] = counts_by_file.get(item.file_id, 0) + 1

    assert response.file_limit == 3
    assert response.file_count == 3
    assert response.has_more is True
    assert counts_by_file
    assert all(count <= 2 for count in counts_by_file.values())


def test_filename_content_search_allows_filename_row_plus_content_limit(tmp_path):
    from backend.api.search import search_files
    from backend.database import save_file_chunks
    from backend.models.schemas import SearchRequest

    named_id = register_file(
        str(tmp_path / '공통키워드-title.docx'),
        '공통키워드-title.docx',
        'Word',
        0,
    )
    save_file_chunks(
        named_id,
        [
            {"location": f"문단 {index}", "content": f"본문 공통키워드 반복 {index}"}
            for index in range(3)
        ],
    )

    other_id = register_file(str(tmp_path / 'other.docx'), 'other.docx', 'Word', 0)
    save_file_chunks(other_id, [{"location": "문단", "content": "본문 공통키워드 다른 파일"}])

    update_file_mtime(named_id, 2.0)
    update_file_mtime(other_id, 1.0)

    response = search_files(
        SearchRequest(query="공통키워드", search_scope="filename_content", file_limit=2, per_file_limit=1)
    )

    named_locations = [item.location for item in response.results if item.file_id == named_id]
    named_content_locations = [location for location in named_locations if location != "파일명"]

    assert response.file_count == 2
    assert "파일명" in named_locations
    assert len(named_content_locations) == 1


def test_content_search_does_not_order_files_by_recent_mtime_and_keeps_document_order(tmp_path):
    first_id = register_file(str(tmp_path / 'old.docx'), 'old.docx', 'Word', 0)
    second_id = register_file(str(tmp_path / 'new.docx'), 'new.docx', 'Word', 0)
    third_id = register_file(str(tmp_path / 'middle.docx'), 'middle.docx', 'Word', 0)
    save_file_chunks(first_id, [{"location": "문단 1", "content": "프로젝트 공통키워드"}])
    save_file_chunks(
        second_id,
        [
            {"location": "슬라이드 1", "content": "프로젝트 공통키워드"},
            {"location": "슬라이드 2", "content": "프로젝트 공통키워드 프로젝트 공통키워드"},
            {"location": "슬라이드 5", "content": "프로젝트 공통키워드"},
        ],
    )
    save_file_chunks(third_id, [{"location": "문단 1", "content": "프로젝트 공통키워드"}])
    update_file_mtime(first_id, 1_000)
    update_file_mtime(second_id, 3_000)
    update_file_mtime(third_id, 2_000)

    results = search("공통키워드", file_limit=3, per_file_limit=3)

    ordered_file_ids = []
    for item in results:
        if item["file_id"] not in ordered_file_ids:
            ordered_file_ids.append(item["file_id"])
    assert ordered_file_ids == [first_id, second_id, third_id]
    assert [item["location"] for item in results if item["file_id"] == second_id] == [
        "슬라이드 1",
        "슬라이드 2",
        "슬라이드 5",
    ]
    assert {item["file_mtime"] for item in results if item["file_id"] == second_id} == {3_000}


def test_search_api_allows_more_files_up_to_max(tmp_path):
    from backend.api.search import search_files
    from backend.database import save_file_chunks
    from backend.models.schemas import SearchRequest

    for index in range(105):
        file_id = register_file(str(tmp_path / f'bulk-{index}.docx'), f'bulk-{index}.docx', 'Word', 0)
        save_file_chunks(file_id, [{"location": "문단", "content": f"프로젝트 대량검색 {index}"}])
        update_file_mtime(file_id, float(index))

    response = search_files(SearchRequest(query="대량검색", search_scope="content", file_limit=120))

    assert response.file_limit == 100
    assert response.file_count == 100
    assert len({item.file_id for item in response.results}) == 100
    assert response.has_more is False


def test_search_api_content_results_include_optional_fingerprint_metadata(tmp_path):
    from backend.api.search import search_files
    from backend.database import save_file_chunks
    from backend.models.schemas import SearchRequest

    file_id = register_file(str(tmp_path / 'fingerprint.docx'), 'fingerprint.docx', 'Word', 0)
    update_file_mtime(file_id, 1_700_000_000)
    save_file_chunks(file_id, [{"location": "문단", "content": "fingerprint body match"}])

    response = search_files(SearchRequest(query="fingerprint", search_scope="content"))

    assert response.total == 1
    result = response.results[0]
    assert result.file_id == file_id
    assert result.normalized_hash
    assert result.content_hash
    assert result.content_chars == len("fingerprint body match")
    assert result.chunk_count == 1
    assert result.file_mtime == 1_700_000_000


def test_search_api_filename_results_include_optional_fingerprint_metadata(tmp_path):
    from backend.api.search import search_files
    from backend.database import save_file_chunks
    from backend.models.schemas import SearchRequest

    file_id = register_file(str(tmp_path / 'filename-hash.docx'), 'filename-hash.docx', 'Word', 0)
    save_file_chunks(file_id, [{"location": "문단", "content": "same content available for filename result"}])

    response = search_files(SearchRequest(query="filename-hash", search_scope="filename"))

    assert response.total == 1
    result = response.results[0]
    assert result.file_id == file_id
    assert result.normalized_hash
    assert result.content_hash
    assert result.content_chars == len("same content available for filename result")
    assert result.chunk_count == 1
