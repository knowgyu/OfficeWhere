from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.util import Inches


def _setup_db(tmp_path, monkeypatch):
    from backend.database import init_db

    monkeypatch.setattr("backend.database.DB_PATH", tmp_path / "test.db")
    monkeypatch.setattr("backend.database.DB_DIR", tmp_path)
    init_db()


def _write_word(path: Path, body: str):
    document = Document()
    document.add_paragraph("공통 서론")
    document.add_paragraph(body)
    document.save(path)


def _write_ppt(path: Path, body: str):
    presentation = Presentation()
    layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = "Overview"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4), Inches(1))
    textbox.text_frame.text = body
    presentation.save(path)


def _index_file(path: Path):
    from backend.core.indexer import inspect_and_chunk
    from backend.database import save_indexed_file

    info, chunks = inspect_and_chunk(str(path))
    return save_indexed_file(
        path=str(path),
        name=info["name"],
        file_type=info["file_type"],
        key_column="",
        column_count=len(info["columns"]),
        chunks=chunks,
        file_mtime=path.stat().st_mtime,
        parser_config=None,
        excel_sheets=info.get("excel_sheets"),
        excel_cells=info.get("excel_cells"),
        comparison_artifacts=info.get("comparison_artifacts"),
    )


def _file_infos(file_ids):
    from backend.database import get_file_by_id

    return [get_file_by_id(file_id) for file_id in file_ids]


def test_word_compare_uses_valid_artifact_without_source_parse(tmp_path, monkeypatch):
    from backend.core.checker import run_consistency_check

    _setup_db(tmp_path, monkeypatch)
    left = tmp_path / "left.docx"
    right = tmp_path / "right.docx"
    _write_word(left, "본문 버전 A")
    _write_word(right, "본문 버전 B")
    file_ids = [_index_file(left), _index_file(right)]

    def fail_source_parse(_path):
        raise AssertionError("valid Word artifact should avoid source parse")

    monkeypatch.setattr("backend.core.word_compare.extract_word_blocks", fail_source_parse)

    result = run_consistency_check(_file_infos(file_ids))

    assert result["mode"] == "word"
    assert any(change["change_type"] == "replace" for change in result["word"]["changes"])
    assert result["metadata"]["artifact_status"] == "ok"
    assert result["metadata"]["warnings"] == []


def test_ppt_compare_uses_valid_artifact_without_source_parse(tmp_path, monkeypatch):
    from backend.core.checker import run_consistency_check

    _setup_db(tmp_path, monkeypatch)
    left = tmp_path / "left.pptx"
    right = tmp_path / "right.pptx"
    _write_ppt(left, "원본 본문")
    _write_ppt(right, "수정 본문")
    file_ids = [_index_file(left), _index_file(right)]

    def fail_source_parse(_path):
        raise AssertionError("valid PPT artifact should avoid source parse")

    monkeypatch.setattr("backend.core.ppt_compare.extract_ppt_slides", fail_source_parse)

    result = run_consistency_check(_file_infos(file_ids))

    assert result["mode"] == "ppt"
    assert any(change["change_type"] == "slide_update" for change in result["ppt"]["changes"])
    assert result["metadata"]["artifact_status"] == "ok"
    assert result["metadata"]["warnings"] == []


def test_artifact_payload_is_compressed_and_size_accounted(tmp_path, monkeypatch):
    from backend.database import PPT_COMPARISON_ARTIFACT_KIND, get_comparison_artifact

    _setup_db(tmp_path, monkeypatch)
    ppt = tmp_path / "deck.pptx"
    _write_ppt(ppt, "본문")
    file_id = _index_file(ppt)

    artifact = get_comparison_artifact(file_id, PPT_COMPARISON_ARTIFACT_KIND)

    assert artifact["status"] == "ok"
    assert artifact["raw_size_bytes"] > 0
    assert artifact["compressed_size_bytes"] > 0
    assert artifact["payload"]["slides"][0]["title"] == "Overview"


def test_artifact_version_mismatch_falls_back_with_warning(tmp_path, monkeypatch):
    from backend.core.checker import run_consistency_check
    from backend.database import WORD_COMPARISON_ARTIFACT_KIND, get_file_by_id, save_comparison_artifact, save_indexed_file

    _setup_db(tmp_path, monkeypatch)
    left = tmp_path / "left.docx"
    right = tmp_path / "right.docx"
    _write_word(left, "본문 버전 A")
    _write_word(right, "본문 버전 B")
    left_id = save_indexed_file(
        path=str(left),
        name=left.name,
        file_type="Word",
        key_column="",
        column_count=0,
        chunks=[],
        file_mtime=left.stat().st_mtime,
    )
    right_id = save_indexed_file(
        path=str(right),
        name=right.name,
        file_type="Word",
        key_column="",
        column_count=0,
        chunks=[],
        file_mtime=right.stat().st_mtime,
    )
    save_comparison_artifact(
        left_id,
        file_type="Word",
        artifact_kind=WORD_COMPARISON_ARTIFACT_KIND,
        artifact_version="old",
        parser_version="word-blocks-v1",
        payload={"blocks": []},
    )

    result = run_consistency_check([get_file_by_id(left_id), get_file_by_id(right_id)])

    assert result["mode"] == "word"
    warning_types = {warning["type"] for warning in result["metadata"]["warnings"]}
    assert "artifact_version_mismatch" in warning_types
    assert "artifact_missing" in warning_types
    assert result["metadata"]["used_last_index_snapshot"] is False


def test_compare_falls_back_when_artifact_database_is_unavailable(tmp_path, monkeypatch):
    from backend.core.ppt_compare import compare_ppt_files
    from backend.core.word_compare import compare_word_files
    from backend import database

    missing_dir = tmp_path / "missing-db-dir"
    monkeypatch.setattr(database, "DB_DIR", missing_dir)
    monkeypatch.setattr(database, "DB_PATH", missing_dir / "data.db")

    left_word = tmp_path / "left.docx"
    right_word = tmp_path / "right.docx"
    _write_word(left_word, "본문 버전 A")
    _write_word(right_word, "본문 버전 B")

    word_result = compare_word_files([
        {"id": 101, "path": str(left_word), "name": left_word.name, "file_type": "Word"},
        {"id": 102, "path": str(right_word), "name": right_word.name, "file_type": "Word"},
    ])

    assert word_result["metadata"]["artifact_status"] == "unavailable"
    assert word_result["metadata"]["used_last_index_snapshot"] is False
    assert any(change["change_type"] == "replace" for change in word_result["changes"])

    left_ppt = tmp_path / "left.pptx"
    right_ppt = tmp_path / "right.pptx"
    _write_ppt(left_ppt, "원본 본문")
    _write_ppt(right_ppt, "수정 본문")

    ppt_result = compare_ppt_files([
        {"id": 201, "path": str(left_ppt), "name": left_ppt.name, "file_type": "PowerPoint"},
        {"id": 202, "path": str(right_ppt), "name": right_ppt.name, "file_type": "PowerPoint"},
    ])

    assert ppt_result["metadata"]["artifact_status"] == "unavailable"
    assert ppt_result["metadata"]["used_last_index_snapshot"] is False
    assert any(change["change_type"] == "slide_update" for change in ppt_result["changes"])
